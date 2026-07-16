"""Inspect the generated PPTX and verify icon placement on each slide."""
import logging
logging.basicConfig(level=logging.DEBUG)

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

prs = Presentation("out/icon_test/output.pptx")
print(f"\nTotal slides: {len(prs.slides)}")
for i, slide in enumerate(prs.slides):
    print(f"\n--- Slide {i+1:02d} ---")
    for s in slide.shapes:
        print(f"  {s.name!r:30s}  shape_type={s.shape_type}  left={s.left/914400:.2f}in  top={s.top/914400:.2f}in")
