"""Quick render test for Step 3 icon placement.

Run from Back-End/ with:
    .venv/Scripts/python.exe render_icon_test.py

Produces out/icon_test/output.pptx and exports slide-02.png / slide-06.png
via LibreOffice headless so you can check icon relevance and position.
"""

import subprocess
import sys
from pathlib import Path

# --- add Back-End to sys.path so app.* imports work ---
HERE = Path(__file__).parent
sys.path.insert(0, str(HERE))

from app.config import settings
from app.models.enums import AssetSlot, AssetType, SlideType
from app.models.job import SlidePlan
from app.services.asset_store import InMemoryAssetStore, init_store
from app.services.composer import compose
from app.services.parser import ParsedDeck, ParsedSlide
from app.services.seeder import seed_all

# ── seed asset store ────────────────────────────────────────────────────────
store = InMemoryAssetStore()
n = seed_all(settings.assets_root, store)
init_store(store)
print(f"[seed] {n} assets loaded")

# ── synthetic slides (iMocha topics, intentionally varied for icon diversity) ─
def _slide(index: int, title: str, bullets: list[str]) -> ParsedSlide:
    return ParsedSlide(
        index=index,
        slide_type=SlideType.content,
        text_blocks=[],
        tables=[],
        images=[],
        claims=[],
        title=title,
        body_items=bullets,
    )


SLIDES = [
    _slide(0, "Welcome to iMocha AI Studio", [
        "iMocha AI Studio transforms your existing PowerPoint presentations into branded iMocha decks automatically.",
        "The system preserves every claim and data point from your original content while applying consistent visual design.",
        "Slide layouts are chosen intelligently based on your content structure — bullet-heavy, data-heavy, or executive summary.",
    ]),
    _slide(1, "AI-Powered Skill Intelligence Platform", [
        "iMocha's AI engine analyses over 2,500 skill taxonomies to map your workforce capabilities against industry benchmarks.",
        "Machine learning algorithms continuously refine skill scores based on assessment performance and market data.",
        "Generative AI assistants provide personalised learning pathways for every employee based on their current skill profile.",
        "Automated skill-gap analysis predicts which competencies will be critical in 6, 12, and 24 months.",
    ]),
    _slide(2, "Talent Acquisition & Hiring Pipeline", [
        "Structured skill assessments at the top of the hiring funnel reduce time-to-hire by 40% on average across enterprise customers.",
        "Candidate scoring is objective and role-specific, eliminating resume bias and ensuring every shortlist is merit-based.",
        "Integration with ATS platforms means assessment data flows directly into recruiter workflows.",
        "Real-time hiring pipeline analytics give talent teams full visibility from application to offer acceptance.",
    ]),
    _slide(3, "Learning & Development at Scale", [
        "Skill-gap reports auto-generate recommended learning paths aligned to each employee's role and career trajectory.",
        "Training content is curated from 50+ LMS providers and surfaced in the employee's daily workflow without context switching.",
        "Completion rates and knowledge retention are tracked through periodic micro-assessments embedded in existing tools.",
        "Upskilling ROI dashboards show L&D leaders which programs are driving measurable performance improvement.",
    ]),
    _slide(4, "Enterprise Team Collaboration", [
        "Shared skill profiles across teams enable managers to identify internal talent for cross-functional projects instantly.",
        "Communication dashboards surface collaboration patterns and flag where knowledge transfer between departments is breaking down.",
        "Team health scores aggregate individual skill trends to help HR business partners spot retention risk early.",
    ]),
    _slide(5, "Analytics & Reporting Dashboard", [
        "Real-time analytics dashboards give CHROs a live view of workforce capability across every business unit and geography.",
        "Benchmark reports compare your organisation's skill maturity against 3,000+ companies in the iMocha data consortium.",
        "Custom KPI tracking lets you measure assessment pass rates, score distributions, and skill coverage against targets.",
        "Automated monthly reports are delivered to stakeholders with zero manual effort from the HR analytics team.",
    ]),
    _slide(6, "Assessment Accuracy & Scoring", [
        "iMocha assessments are validated by a team of 200+ subject-matter experts across 15 domains to ensure high accuracy.",
        "Anti-cheating AI monitors behavioural patterns during assessments to maintain score integrity at scale.",
        "Scoring rubrics are calibrated against industry certification benchmarks so results are immediately interpretable.",
        "Pass/fail thresholds are configurable per role so hiring managers can tailor standards to their specific needs.",
    ]),
]

plan = [
    SlidePlan(
        id=f"slide-{s.index:02d}",
        index=s.index,
        slide_type=s.slide_type,
        planned_layout="body-block",
        asset_types=[],
        layout_category="body-block",
    )
    for s in SLIDES
]

parsed = ParsedDeck(name="icon_test", slide_count=len(SLIDES), slides=SLIDES)

# ── compose ──────────────────────────────────────────────────────────────────
print("[compose] building deck ...")
prs, flagged = compose(parsed, plan, settings.assets_root)
if flagged:
    print(f"[compose] overflow flagged on slide indices: {flagged}")

out_dir = HERE / "out" / "icon_test"
out_dir.mkdir(parents=True, exist_ok=True)
out_pptx = out_dir / "output.pptx"
prs.save(str(out_pptx))
print(f"[compose] saved -> {out_pptx}")

# ── inspect slides 2 and 6 ───────────────────────────────────────────────────
from pptx import Presentation as _Prs
from pptx.util import Emu as _Emu

_check = _Prs(str(out_pptx))
for slide_1based in (2, 6):
    slide = _check.slides[slide_1based - 1]
    pics = [s for s in slide.shapes if s.shape_type == 13]  # PICTURE
    # Body textbox: find the largest textbox (by area) — it's the injected body
    txts = [s for s in slide.shapes if s.shape_type == 17]  # TEXT_BOX
    body_box = max(txts, key=lambda s: s.width * s.height, default=None)
    font_pt = None
    if body_box:
        tf = body_box.text_frame
        if tf.paragraphs and tf.paragraphs[0].runs:
            font_pt = tf.paragraphs[0].runs[0].font.size
            font_pt = font_pt / 12700 if font_pt else None  # EMU -> pt
    # Icon: the last picture (logo is Picture 2 at y~6.88in; icon is the other one)
    icon_pic = next(
        (p for p in pics if p.top / 914400 < 3.0),  # icon is in top half
        None,
    )
    icon_name = "none"
    icon_in = 0.0
    if icon_pic:
        icon_in = icon_pic.width / 914400
        # Retrieve filename from the relationship
        try:
            rId = icon_pic.element.find(
                ".//{http://schemas.openxmlformats.org/drawingml/2006/spreadsheetDrawing}blipFill",
            )
            # Simpler: just report dimensions
        except Exception:
            pass
        icon_name = f"w={icon_in:.2f}in  h={icon_pic.height/914400:.2f}in  left={icon_pic.left/914400:.2f}in  top={icon_pic.top/914400:.2f}in"
    font_str = f"{font_pt:.1f}pt" if font_pt else "(not found)"
    print(f"\n[inspect] Slide {slide_1based:02d}: '{SLIDES[slide_1based-1].title}'")
    print(f"  body font size : {font_str}")
    print(f"  icon shape     : {icon_name}")

# ── export slides 02 and 06 via LibreOffice headless ────────────────────────
lo_candidates = [
    Path(r"C:\Program Files\LibreOffice\program\soffice.exe"),
    Path(r"C:\Program Files (x86)\LibreOffice\program\soffice.exe"),
]
lo = next((p for p in lo_candidates if p.exists()), None)
if lo is None:
    print("[export] LibreOffice not found — skipping PNG export")
    sys.exit(0)

png_dir = out_dir / "png"
png_dir.mkdir(exist_ok=True)

print(f"[export] running LibreOffice headless -> {png_dir}")
result = subprocess.run(
    [
        str(lo),
        "--headless",
        "--convert-to", "png",
        "--outdir", str(png_dir),
        str(out_pptx),
    ],
    capture_output=True, text=True, timeout=120,
)
if result.returncode != 0:
    print("[export] LibreOffice error:", result.stderr[:500])
    sys.exit(1)

print("[export] done. PNG files:")
for f in sorted(png_dir.glob("*.png")):
    print(f"  {f.name}")

# LibreOffice names slides: output1.png, output2.png, …
slide02 = png_dir / "output2.png"
slide06 = png_dir / "output6.png"
for path, label in [(slide02, "slide-02"), (slide06, "slide-06")]:
    if path.exists():
        dest = out_dir / f"{label}.png"
        path.rename(dest)
        print(f"[export] {label} -> {dest}")
    else:
        print(f"[export] WARNING: {path.name} not found")
