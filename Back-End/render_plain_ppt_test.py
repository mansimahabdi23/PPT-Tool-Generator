"""Batch 1 correctness test — cover clone, guide-box removal, closing append.

Run from Back-End/ with:
    .venv/Scripts/python.exe render_plain_ppt_test.py

Produces out/plain_ppt_test/output.pptx and exports PNGs via LibreOffice.

Slides rendered
---------------
Slide 1  : cover  — clone of template slide 1, filled with source title + bullets
Slides 2-4: body-block — NO lavender guide box; icon accent in upper-right
Slide 5  : Bookmark Card infographic — 5 parallel AI capabilities
Slide 6  : Hub-and-Spoke infographic — 4 governance pillars
Slide 7  : last content slide — full content preserved as body-block (NOT closing)
Slide 8  : closing — clone of template slide 20, appended

Content-diff check
------------------
Prints every claim from the source first slide and verifies it appears in the
cover text; prints every claim from the source last slide and verifies it is
present in its own body-block output slide.
"""

import subprocess
import sys
from pathlib import Path

HERE = Path(__file__).parent
sys.path.insert(0, str(HERE))

from app.config import settings
from app.models.enums import AssetSlot, AssetType, SlideType
from app.models.job import SlidePlan
from app.services.asset_store import InMemoryAssetStore, init_store
from app.services.composer import compose
from app.services.parser import ParsedDeck, ParsedSlide
from app.services.seeder import seed_all

# ── seed ────────────────────────────────────────────────────────────────────
store = InMemoryAssetStore()
n = seed_all(settings.assets_root, store)
init_store(store)
print(f"[seed] {n} assets loaded")

# ── source deck (Plain_PPT equivalent — 6 content slides) ───────────────────
# Index 0 = cover (SlideType.title assigned by parser; here we give it title+bullets)
# Index 5 = last source slide (the parser would assign SlideType.closing — we
#           must render it as body-block and keep all bullets)

SOURCE_COVER_TITLE  = "iMocha AI Assessment Platform"
SOURCE_COVER_BULLETS = [
    "Q4 2025 Business Review",
    "The world's most accurate skills intelligence platform — powering over 500 enterprise talent teams.",
    "Covering 2,500+ skill taxonomies across engineering, analytics, leadership, and sales.",
    "Trusted by Ericsson, Capgemini, Cognizant, and 100+ Fortune 500 companies.",
]

SOURCE_LAST_TITLE   = "Next Steps & Call to Action"
SOURCE_LAST_BULLETS = [
    "Schedule a live demo with your iMocha account executive within the next 14 days.",
    "Pilot iMocha on 3 open roles to benchmark assessment accuracy against your current process.",
    "Connect your ATS (Greenhouse, Lever, Workday) via our pre-built integration connectors.",
    "Join the iMocha Skills Community for best-practice sharing with 1,200+ talent leaders.",
]


def _slide(
    index: int,
    slide_type: SlideType,
    title: str,
    bullets: list[str],
) -> ParsedSlide:
    return ParsedSlide(
        index=index,
        slide_type=slide_type,
        text_blocks=[],
        tables=[],
        images=[],
        claims=[],
        title=title,
        body_items=bullets,
    )


SLIDES = [
    _slide(0, SlideType.title, SOURCE_COVER_TITLE, SOURCE_COVER_BULLETS),
    _slide(1, SlideType.content, "AI-Powered Skill Intelligence", [
        "Maps 2,500+ skill taxonomies against live market benchmarks automatically.",
        "Machine-learning scoring adapts to role-specific pass/fail thresholds.",
        "Generative AI assistant builds personalised upskilling paths per employee.",
    ]),
    _slide(2, SlideType.content, "Hiring Pipeline Efficiency", [
        "Structured pre-screen assessments cut time-to-hire by 40% on average.",
        "Objective candidate scoring eliminates resume bias at the top of funnel.",
        "Real-time pipeline analytics: live view from application to offer acceptance.",
    ]),
    _slide(3, SlideType.content, "Learning & Development at Scale", [
        "Auto-generated skill-gap reports align learning paths to each employee's role.",
        "Content curated from 50+ LMS providers, surfaced in the employee's daily workflow.",
        "Upskilling ROI dashboards show which programs drive measurable performance gains.",
    ]),
    _slide(4, SlideType.content, "Core AI Capabilities", [
        "AI Skills Inference Engine",
        "Conversational AI Career Assistant",
        "AI-Powered Job Profile Creation",
        "Intelligent Skill Gap Analysis",
        "Personalized Learning Recommendations",
    ]),
    _slide(5, SlideType.content, "AI Governance, Trust & Transparency", [
        "Transparent and explainable AI-driven matching logic",
        "Bias-controlled workforce recommendations and assessments",
        "Enterprise-grade privacy, security, and compliance architecture",
        "Audit-ready governance workflows for enterprise trust",
    ]),
    # Last source slide — parser would tag this SlideType.closing; our fix keeps it as body-block
    _slide(6, SlideType.closing, SOURCE_LAST_TITLE, SOURCE_LAST_BULLETS),
]

# All slides are body-block in the plan except slide 0 (cover) which has no layout_category
# so it falls through to the SlideType.title clone path in the composer.
plan = [
    SlidePlan(
        id="cover",
        index=0,
        slide_type=SlideType.title,
        planned_layout="cover",
        asset_types=[],
        layout_category=None,   # triggers clone_cover_slide path
    ),
] + [
    SlidePlan(
        id=f"slide-{s.index:02d}",
        index=s.index,
        slide_type=s.slide_type,
        planned_layout="body-block",
        asset_types=[],
        layout_category="body-block",
    )
    for s in SLIDES[1:-1]  # content slides
] + [
    SlidePlan(
        id="slide-last",
        index=6,
        slide_type=SlideType.closing,   # parser-assigned type; composer renders as body-block
        planned_layout="body-block",
        asset_types=[],
        layout_category=None,           # no layout_category → closing intercept path
    ),
]

parsed = ParsedDeck(name="plain_ppt_test", slide_count=len(SLIDES), slides=SLIDES)

# ── compose ──────────────────────────────────────────────────────────────────
print("[compose] building deck ...")
prs, flagged, infographic_indices = compose(parsed, plan, settings.assets_root)
if flagged:
    print(f"[compose] overflow flagged on slide indices: {flagged}")

expected_slides = len(SLIDES) + 1  # +1 for appended closing
actual_slides   = len(prs.slides)
print(f"[compose] slides: expected={expected_slides}  actual={actual_slides}  {'OK' if actual_slides == expected_slides else 'MISMATCH'}")

out_dir = HERE / "out" / "plain_ppt_test"
out_dir.mkdir(parents=True, exist_ok=True)
out_pptx = out_dir / "output.pptx"
prs.save(str(out_pptx))
print(f"[compose] saved -> {out_pptx}")

# ── content-diff check ───────────────────────────────────────────────────────
from pptx import Presentation as _Prs

_check = _Prs(str(out_pptx))

def _slide_text(slide_1based: int) -> str:
    slide = _check.slides[slide_1based - 1]
    parts = []
    for shape in slide.shapes:
        if hasattr(shape, "text_frame"):
            parts.append(shape.text_frame.text)
    return "\n".join(parts)

print("\n--- CONTENT-DIFF: cover (slide 1) ---")
cover_text = _slide_text(1)
all_ok = True
for claim in [SOURCE_COVER_TITLE] + SOURCE_COVER_BULLETS:
    found = claim in cover_text
    status = "OK " if found else "MISSING"
    if not found:
        all_ok = False
    print(f"  [{status}] {claim[:80]}")
print("  -> cover content-diff:", "PASS" if all_ok else "FAIL")

print(f"\n--- CONTENT-DIFF: last content slide (slide {len(SLIDES)}) ---")
last_text = _slide_text(len(SLIDES))
all_ok = True
for claim in [SOURCE_LAST_TITLE] + SOURCE_LAST_BULLETS:
    found = claim in last_text
    status = "OK " if found else "MISSING"
    if not found:
        all_ok = False
    print(f"  [{status}] {claim[:80]}")
print("  -> last-slide content-diff:", "PASS" if all_ok else "FAIL")

print(f"\n--- CLOSING slide (slide {len(prs.slides)}) ---")
closing_text = _slide_text(len(prs.slides))
# Closing is template slide 20 — check it has the brand text
brand_text_found = "Skills Visibility" in closing_text or "Business Agility" in closing_text
print(f"  brand text present: {'YES' if brand_text_found else 'NO'}")
print("  closing check:", "PASS" if brand_text_found else "FAIL")

# ── shape inspection ──────────────────────────────────────────────────────────
print("\n--- SHAPE NAMES on slide 2 (first body-block) ---")
slide2 = _check.slides[1]
for s in slide2.shapes:
    print(f"  {s.name!r:35s}  shape_type={s.shape_type}")

# ── render per-slide PNGs via LibreOffice -> PDF -> pdftoppm ─────────────────
from app.services.exporter import render_previews

print("\n[export] rendering per-slide PNGs via render_previews() ...")
pngs = render_previews("plain_ppt_test", out_pptx, "transformed")
if not pngs:
    print("[export] no PNGs rendered (LibreOffice or pdftoppm unavailable) — skipping")
    sys.exit(0)

print(f"[export] {len(pngs)} PNG(s) rendered:")
for p in pngs:
    print(f"  {p.name}")

import shutil
total = len(prs.slides)
for slide_1based, label in [
    (1, "slide-01-cover"),
    (5, "slide-05-bookmark-infographic"),
    (6, "slide-06-hubspoke-infographic"),
    (total - 1, f"slide-{total-1:02d}-last-content"),
    (total, f"slide-{total:02d}-closing"),
]:
    src = pngs[slide_1based - 1] if slide_1based <= len(pngs) else None
    if src and src.exists():
        dst = out_dir / f"{label}.png"
        shutil.copy(src, dst)
        print(f"[export] {label} -> {dst}")
    else:
        print(f"[export] WARNING: slide-{slide_1based}.png not found")
