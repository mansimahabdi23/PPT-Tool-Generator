"""Tests for the three deterministic validators: brand_lint, content_diff, layout_qa.

Each validator is a pure function — all tests build in-memory Presentations and
ParsedDecks, assert on structured result objects. No disk I/O beyond tmp files.
"""

from __future__ import annotations

import io
import struct
import tempfile
import zlib
from dataclasses import dataclass, field
from pathlib import Path

import pytest
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Emu, Pt
from lxml import etree

from app import brand
from app.models.enums import SlideType
from app.services.brand_lint import BrandLintResult, lint
from app.services.content_diff import ContentDiffResult, diff
from app.services.layout_qa import LayoutQAResult, check
from app.services.parser import (
    AtomicClaim,
    ParsedDeck,
    ParsedSlide,
    ParsedTextBlock,
    ParsedTable,
    ParsedImage,
    ShapeGeometry,
    SourceRef,
)

# ---------------------------------------------------------------------------
# Shared helpers
# ---------------------------------------------------------------------------

_P = "http://schemas.openxmlformats.org/presentationml/2006/main"
_A = "http://schemas.openxmlformats.org/drawingml/2006/main"

_SLIDE_W = brand.SLIDE_WIDTH_EMU
_SLIDE_H = brand.SLIDE_HEIGHT_EMU


def _blank_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width = Emu(_SLIDE_W)
    prs.slide_height = Emu(_SLIDE_H)
    return prs


def _add_blank_slide(prs: Presentation):
    layout = prs.slide_layouts[6]  # blank
    return prs.slides.add_slide(layout)


def _add_text_run(
    slide,
    text: str,
    font_name: str | None = None,
    color_hex: str | None = None,
    left: int = 500_000,
    top: int = 500_000,
    width: int = 3_000_000,
    height: int = 500_000,
):
    """Add a text box with a single run, optionally styled."""
    box = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))
    p = box.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = text
    if font_name:
        run.font.name = font_name
    if color_hex:
        r, g, b = brand.rgb_tuple(color_hex)
        run.font.color.rgb = RGBColor(r, g, b)
    return box


def _add_solid_fill_shape(slide, color_hex: str, left: int = 500_000, top: int = 500_000):
    """Add a rectangle with an explicit solid fill color."""
    shape = slide.shapes.add_shape(1, Emu(left), Emu(top), Emu(500_000), Emu(300_000))
    shape.fill.solid()
    r, g, b = brand.rgb_tuple(color_hex)
    shape.fill.fore_color.rgb = RGBColor(r, g, b)
    return shape


def _tiny_png() -> bytes:
    """Return a minimal valid 1×1 white PNG."""
    def chunk(ctype: bytes, data: bytes) -> bytes:
        c = zlib.crc32(ctype + data) & 0xFFFFFFFF
        return struct.pack(">I", len(data)) + ctype + data + struct.pack(">I", c)

    header = b"\x89PNG\r\n\x1a\n"
    ihdr = chunk(b"IHDR", struct.pack(">IIBBBBB", 1, 1, 8, 2, 0, 0, 0))
    raw = b"\x00\xFF\xFF\xFF"
    idat = chunk(b"IDAT", zlib.compress(raw))
    iend = chunk(b"IEND", b"")
    return header + ihdr + idat + iend


def _add_logo_picture(slide, top: int | None = None):
    """Add a small picture in the footer region (top > 80% of slide height)."""
    if top is None:
        top = int(_SLIDE_H * 0.85)
    img_bytes = _tiny_png()
    with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as f:
        f.write(img_bytes)
        tmp_path = f.name
    slide.shapes.add_picture(tmp_path, Emu(100_000), Emu(top), width=Emu(400_000), height=Emu(80_000))


def _set_slide_gradient_bg(slide, stop1: str, stop2: str) -> None:
    """Set slide background to a two-stop gradient (replicates builder logic)."""
    cSld = slide.element.find(f"{{{_P}}}cSld")
    for existing in cSld.findall(f"{{{_P}}}bg"):
        cSld.remove(existing)
    bg_xml = (
        f'<p:bg xmlns:p="{_P}" xmlns:a="{_A}">'
        f"  <p:bgPr>"
        f'    <a:gradFill rot="1">'
        f"      <a:gsLst>"
        f'        <a:gs pos="0"><a:srgbClr val="{stop1}"/></a:gs>'
        f'        <a:gs pos="100000"><a:srgbClr val="{stop2}"/></a:gs>'
        f"      </a:gsLst>"
        f"    </a:gradFill>"
        f"  </p:bgPr>"
        f"</p:bg>"
    )
    cSld.insert(0, etree.fromstring(bg_xml))


def _empty_geom() -> ShapeGeometry:
    return ShapeGeometry(0, 0, 0, 0, 0.0, 0.0, 0.0, 0.0)


def _make_source_ref(slide_index: int = 0, shape_id: int = 1) -> SourceRef:
    return SourceRef(slide_index=slide_index, shape_id=shape_id, shape_name="test_shape")


def _make_claim(text: str, ctype: str = "bullet", slide_index: int = 0) -> AtomicClaim:
    return AtomicClaim(text=text, claim_type=ctype, source=_make_source_ref(slide_index))


def _make_empty_slide_model(index: int = 0) -> ParsedSlide:
    return ParsedSlide(
        index=index,
        slide_type=SlideType.content,
        text_blocks=[],
        tables=[],
        images=[],
        claims=[],
        title="",
        body_items=[],
    )


def _parsed_deck(claims: list[AtomicClaim]) -> ParsedDeck:
    slide = _make_empty_slide_model()
    slide.claims = claims
    return ParsedDeck(name="test.pptx", slide_count=1, slides=[slide])


# ---------------------------------------------------------------------------
# 1. brand_lint — fonts
# ---------------------------------------------------------------------------


class TestBrandLintFonts:
    def test_allowed_font_no_violation(self) -> None:
        prs = _blank_prs()
        slide = _add_blank_slide(prs)
        _add_text_run(slide, "Hello", font_name=brand.PLAYFAIR)
        _add_logo_picture(slide)
        result = lint(prs)
        font_violations = [v for v in result.violations if v.violation_type == "font"]
        assert font_violations == []

    def test_poppins_no_violation(self) -> None:
        prs = _blank_prs()
        slide = _add_blank_slide(prs)
        _add_text_run(slide, "Body text", font_name=brand.POPPINS)
        _add_logo_picture(slide)
        result = lint(prs)
        assert not any(v.violation_type == "font" for v in result.violations)

    def test_off_brand_font_flagged(self) -> None:
        prs = _blank_prs()
        slide = _add_blank_slide(prs)
        _add_text_run(slide, "Bad font", font_name="Arial")
        _add_logo_picture(slide)
        result = lint(prs)
        font_violations = [v for v in result.violations if v.violation_type == "font"]
        assert len(font_violations) == 1
        assert "Arial" in font_violations[0].detail

    def test_multiple_off_brand_fonts_each_flagged(self) -> None:
        prs = _blank_prs()
        slide = _add_blank_slide(prs)
        _add_text_run(slide, "Times", font_name="Times New Roman")
        _add_text_run(slide, "Comic", font_name="Comic Sans MS", top=700_000)
        _add_logo_picture(slide)
        result = lint(prs)
        font_violations = [v for v in result.violations if v.violation_type == "font"]
        assert len(font_violations) == 2

    def test_no_font_set_no_violation(self) -> None:
        prs = _blank_prs()
        slide = _add_blank_slide(prs)
        _add_text_run(slide, "Inherited font")  # font_name=None
        _add_logo_picture(slide)
        result = lint(prs)
        assert not any(v.violation_type == "font" for v in result.violations)


# ---------------------------------------------------------------------------
# 2. brand_lint — colors
# ---------------------------------------------------------------------------


class TestBrandLintColors:
    def test_brand_orange_no_violation(self) -> None:
        prs = _blank_prs()
        slide = _add_blank_slide(prs)
        _add_text_run(slide, "Title", color_hex=brand.BRAND_ORANGE)
        _add_logo_picture(slide)
        result = lint(prs)
        assert not any(v.violation_type == "color" for v in result.violations)

    def test_ink_color_no_violation(self) -> None:
        prs = _blank_prs()
        slide = _add_blank_slide(prs)
        _add_text_run(slide, "Body", color_hex=brand.INK)
        _add_logo_picture(slide)
        result = lint(prs)
        assert not any(v.violation_type == "color" for v in result.violations)

    def test_off_brand_text_color_flagged(self) -> None:
        prs = _blank_prs()
        slide = _add_blank_slide(prs)
        _add_text_run(slide, "Wrong color", color_hex="FF0000")
        _add_logo_picture(slide)
        result = lint(prs)
        color_violations = [v for v in result.violations if v.violation_type == "color"]
        assert len(color_violations) >= 1
        assert "FF0000" in color_violations[0].detail

    def test_allowed_shape_fill_no_violation(self) -> None:
        prs = _blank_prs()
        slide = _add_blank_slide(prs)
        _add_solid_fill_shape(slide, brand.BRAND_ORANGE)
        _add_logo_picture(slide)
        result = lint(prs)
        assert not any(v.violation_type == "color" for v in result.violations)

    def test_off_brand_shape_fill_flagged(self) -> None:
        prs = _blank_prs()
        slide = _add_blank_slide(prs)
        _add_solid_fill_shape(slide, "00FF00")  # bright green — not allowed
        _add_logo_picture(slide)
        result = lint(prs)
        color_violations = [v for v in result.violations if v.violation_type == "color"]
        assert any("00FF00" in v.detail for v in color_violations)


# ---------------------------------------------------------------------------
# 3. brand_lint — gradients
# ---------------------------------------------------------------------------


class TestBrandLintGradients:
    def test_allowed_bg_gradient_no_violation(self) -> None:
        prs = _blank_prs()
        slide = _add_blank_slide(prs)
        _set_slide_gradient_bg(slide, brand.BRAND_ORANGE, brand.BRAND_PURPLE)
        _add_logo_picture(slide)
        result = lint(prs)
        assert not any(v.violation_type == "gradient" for v in result.violations)

    def test_purple_indigo_gradient_no_violation(self) -> None:
        prs = _blank_prs()
        slide = _add_blank_slide(prs)
        _set_slide_gradient_bg(slide, brand.BRAND_PURPLE, brand.BRAND_INDIGO)
        _add_logo_picture(slide)
        result = lint(prs)
        assert not any(v.violation_type == "gradient" for v in result.violations)

    def test_off_brand_bg_gradient_flagged(self) -> None:
        prs = _blank_prs()
        slide = _add_blank_slide(prs)
        _set_slide_gradient_bg(slide, "FF0000", "00FF00")  # red→green: not allowed
        _add_logo_picture(slide)
        result = lint(prs)
        grad_violations = [v for v in result.violations if v.violation_type == "gradient"]
        assert len(grad_violations) == 1
        assert "FF0000" in grad_violations[0].detail

    def test_no_gradient_bg_no_gradient_violation(self) -> None:
        prs = _blank_prs()
        slide = _add_blank_slide(prs)
        _add_logo_picture(slide)
        result = lint(prs)
        assert not any(v.violation_type == "gradient" for v in result.violations)


# ---------------------------------------------------------------------------
# 4. brand_lint — logo
# ---------------------------------------------------------------------------


class TestBrandLintLogo:
    def test_logo_in_footer_passes(self) -> None:
        prs = _blank_prs()
        slide = _add_blank_slide(prs)
        _add_logo_picture(slide, top=int(_SLIDE_H * 0.85))
        result = lint(prs)
        assert not any(v.violation_type == "logo_missing" for v in result.violations)

    def test_no_picture_flags_logo_missing(self) -> None:
        prs = _blank_prs()
        _add_blank_slide(prs)
        result = lint(prs)
        logo_violations = [v for v in result.violations if v.violation_type == "logo_missing"]
        assert len(logo_violations) == 1

    def test_picture_named_logo_passes(self) -> None:
        prs = _blank_prs()
        slide = _add_blank_slide(prs)
        _add_logo_picture(slide, top=int(_SLIDE_H * 0.85))
        # Rename it to include "logo"
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                shape.name = "logo_footer"
                break
        result = lint(prs)
        assert not any(v.violation_type == "logo_missing" for v in result.violations)

    def test_picture_above_footer_threshold_flags_missing(self) -> None:
        """A picture placed in the top half of the slide is NOT considered a logo."""
        prs = _blank_prs()
        slide = _add_blank_slide(prs)
        # top = 10% of slide height — well above the 80% threshold
        _add_logo_picture(slide, top=int(_SLIDE_H * 0.10))
        result = lint(prs)
        logo_violations = [v for v in result.violations if v.violation_type == "logo_missing"]
        assert len(logo_violations) == 1

    def test_overall_passed_true_when_no_violations(self) -> None:
        prs = _blank_prs()
        slide = _add_blank_slide(prs)
        _add_text_run(slide, "Title", font_name=brand.PLAYFAIR, color_hex=brand.BRAND_ORANGE)
        _add_logo_picture(slide)
        result = lint(prs)
        assert result.passed is True
        assert result.violations == []


# ---------------------------------------------------------------------------
# 5. brand_lint — shapes_checked / slide_count metadata
# ---------------------------------------------------------------------------


class TestBrandLintMetadata:
    def test_shapes_checked_counts_text_boxes(self) -> None:
        prs = _blank_prs()
        slide = _add_blank_slide(prs)
        _add_text_run(slide, "A")
        _add_text_run(slide, "B", top=700_000)
        _add_logo_picture(slide)
        result = lint(prs)
        assert result.shapes_checked >= 2  # at least the two text boxes

    def test_slide_count_matches(self) -> None:
        prs = _blank_prs()
        for _ in range(3):
            slide = _add_blank_slide(prs)
            _add_logo_picture(slide)
        result = lint(prs)
        assert result.slide_count == 3


# ---------------------------------------------------------------------------
# 6. content_diff — basic
# ---------------------------------------------------------------------------


class TestContentDiffBasic:
    def test_empty_source_passes(self) -> None:
        prs = _blank_prs()
        deck = ParsedDeck(name="empty.pptx", slide_count=0, slides=[])
        result = diff(deck, prs)
        assert result.passed is True
        assert result.total_claims == 0
        assert result.preserved_claims == 0
        assert "0/0" in result.fidelity_str

    def test_claim_present_passes(self) -> None:
        prs = _blank_prs()
        slide = _add_blank_slide(prs)
        _add_text_run(slide, "Hiring quality improved 34%")
        deck = _parsed_deck([_make_claim("Hiring quality improved 34%", "statistic")])
        result = diff(deck, prs)
        assert result.passed is True
        assert result.total_claims == 1
        assert result.preserved_claims == 1
        assert result.missing_claims == []

    def test_claim_absent_fails(self) -> None:
        prs = _blank_prs()
        _add_blank_slide(prs)  # no text
        deck = _parsed_deck([_make_claim("Critical missing fact")])
        result = diff(deck, prs)
        assert result.passed is False
        assert result.total_claims == 1
        assert result.preserved_claims == 0
        assert len(result.missing_claims) == 1
        assert result.missing_claims[0].text == "Critical missing fact"

    def test_partial_preservation(self) -> None:
        prs = _blank_prs()
        slide = _add_blank_slide(prs)
        _add_text_run(slide, "Present claim")
        deck = _parsed_deck([
            _make_claim("Present claim"),
            _make_claim("Missing claim"),
        ])
        result = diff(deck, prs)
        assert result.passed is False
        assert result.total_claims == 2
        assert result.preserved_claims == 1
        assert "1/2" in result.fidelity_str

    def test_all_preserved_fidelity_string(self) -> None:
        prs = _blank_prs()
        slide = _add_blank_slide(prs)
        _add_text_run(slide, "Claim A")
        _add_text_run(slide, "Claim B", top=700_000)
        deck = _parsed_deck([
            _make_claim("Claim A"),
            _make_claim("Claim B"),
        ])
        result = diff(deck, prs)
        assert result.passed is True
        assert "2/2" in result.fidelity_str


# ---------------------------------------------------------------------------
# 7. content_diff — normalization
# ---------------------------------------------------------------------------


class TestContentDiffNormalization:
    def test_case_insensitive_match(self) -> None:
        prs = _blank_prs()
        slide = _add_blank_slide(prs)
        _add_text_run(slide, "HIRING QUALITY IMPROVED")
        deck = _parsed_deck([_make_claim("hiring quality improved")])
        result = diff(deck, prs)
        assert result.passed is True

    def test_whitespace_normalization(self) -> None:
        prs = _blank_prs()
        slide = _add_blank_slide(prs)
        _add_text_run(slide, "Hiring   quality\nimproved")
        deck = _parsed_deck([_make_claim("Hiring quality improved")])
        result = diff(deck, prs)
        assert result.passed is True

    def test_statistic_exact_number_preserved(self) -> None:
        """34% must appear exactly — not just 3 or 4%."""
        prs = _blank_prs()
        slide = _add_blank_slide(prs)
        _add_text_run(slide, "Accuracy increased by 34% year over year")
        deck = _parsed_deck([_make_claim("increased by 34% year over year", "statistic")])
        result = diff(deck, prs)
        assert result.passed is True

    def test_wrong_number_fails(self) -> None:
        prs = _blank_prs()
        slide = _add_blank_slide(prs)
        _add_text_run(slide, "Accuracy increased by 35% year over year")
        deck = _parsed_deck([_make_claim("increased by 34%", "statistic")])
        result = diff(deck, prs)
        assert result.passed is False

    def test_table_cell_text_found_in_output(self) -> None:
        """Content extracted from output tables must satisfy claims."""
        prs = _blank_prs()
        slide = _add_blank_slide(prs)
        # Add a simple 1×1 table
        table_shape = slide.shapes.add_table(1, 1, Emu(1_000_000), Emu(1_000_000),
                                              Emu(2_000_000), Emu(500_000)).table
        table_shape.cell(0, 0).text_frame.paragraphs[0].runs  # ensure writable
        from pptx.util import Pt as _Pt
        tf = table_shape.cell(0, 0).text_frame
        tf.paragraphs[0].text = "Revenue Q4"
        deck = _parsed_deck([_make_claim("Revenue Q4", "table_cell")])
        result = diff(deck, prs)
        assert result.passed is True

    def test_missing_claim_has_correct_metadata(self) -> None:
        prs = _blank_prs()
        _add_blank_slide(prs)
        claim = AtomicClaim(
            text="Specific stat 99%",
            claim_type="statistic",
            source=SourceRef(slide_index=2, shape_id=5, shape_name="data_box"),
        )
        deck = _parsed_deck([claim])
        result = diff(deck, prs)
        assert result.missing_claims[0].claim_type == "statistic"
        assert result.missing_claims[0].source_slide == 2
        assert result.missing_claims[0].source_shape == "data_box"


# ---------------------------------------------------------------------------
# 8. layout_qa — overflow
# ---------------------------------------------------------------------------


class TestLayoutQAOverflow:
    def test_shape_within_bounds_no_issue(self) -> None:
        prs = _blank_prs()
        slide = _add_blank_slide(prs)
        # Shape entirely within the slide
        slide.shapes.add_shape(1, Emu(1_000_000), Emu(500_000), Emu(2_000_000), Emu(1_000_000))
        result = check(prs)
        assert not any(i.issue_type == "overflow" for i in result.issues)

    def test_shape_overflow_right(self) -> None:
        prs = _blank_prs()
        slide = _add_blank_slide(prs)
        # left=11M + width=2M = 13M > slide_width=12.19M
        slide.shapes.add_shape(1, Emu(11_000_000), Emu(500_000), Emu(2_000_000), Emu(500_000))
        result = check(prs)
        overflow = [i for i in result.issues if i.issue_type == "overflow"]
        assert len(overflow) == 1
        assert "right" in overflow[0].detail

    def test_shape_overflow_bottom(self) -> None:
        prs = _blank_prs()
        slide = _add_blank_slide(prs)
        # top=6M + height=1.5M = 7.5M > slide_height=6.858M
        slide.shapes.add_shape(1, Emu(500_000), Emu(6_000_000), Emu(1_000_000), Emu(1_500_000))
        result = check(prs)
        overflow = [i for i in result.issues if i.issue_type == "overflow"]
        assert len(overflow) == 1
        assert "bottom" in overflow[0].detail

    def test_shape_exactly_at_boundary_no_overflow(self) -> None:
        """A shape whose right edge exactly equals the slide width must not flag."""
        prs = _blank_prs()
        slide = _add_blank_slide(prs)
        # right = slide_width exactly
        slide.shapes.add_shape(1, Emu(0), Emu(0), Emu(_SLIDE_W), Emu(_SLIDE_H // 2))
        result = check(prs)
        assert not any(i.issue_type == "overflow" for i in result.issues)


# ---------------------------------------------------------------------------
# 9. layout_qa — off-screen
# ---------------------------------------------------------------------------


class TestLayoutQAOffScreen:
    def test_negative_left_flagged(self) -> None:
        prs = _blank_prs()
        slide = _add_blank_slide(prs)
        shape = slide.shapes.add_shape(1, Emu(500_000), Emu(500_000), Emu(1_000_000), Emu(500_000))
        # Force negative left via element manipulation
        shape._element.spPr.xfrm.off.set("x", str(-100_000))
        result = check(prs)
        off_screen = [i for i in result.issues if i.issue_type == "off_screen"]
        assert len(off_screen) == 1

    def test_zero_left_zero_top_ok(self) -> None:
        prs = _blank_prs()
        slide = _add_blank_slide(prs)
        slide.shapes.add_shape(1, Emu(0), Emu(0), Emu(1_000_000), Emu(500_000))
        result = check(prs)
        assert not any(i.issue_type == "off_screen" for i in result.issues)


# ---------------------------------------------------------------------------
# 10. layout_qa — empty slide
# ---------------------------------------------------------------------------


class TestLayoutQAEmptySlide:
    def test_slide_with_text_not_empty(self) -> None:
        prs = _blank_prs()
        slide = _add_blank_slide(prs)
        _add_text_run(slide, "Hello world")
        result = check(prs)
        assert not any(i.issue_type == "empty_slide" for i in result.issues)

    def test_slide_with_no_content_flagged(self) -> None:
        prs = _blank_prs()
        _add_blank_slide(prs)  # truly empty
        result = check(prs)
        empty = [i for i in result.issues if i.issue_type == "empty_slide"]
        assert len(empty) == 1

    def test_slide_with_only_whitespace_text_flagged(self) -> None:
        prs = _blank_prs()
        slide = _add_blank_slide(prs)
        _add_text_run(slide, "   \n  ")  # whitespace only
        result = check(prs)
        empty = [i for i in result.issues if i.issue_type == "empty_slide"]
        assert len(empty) == 1

    def test_slide_with_full_bleed_image_only_flagged(self) -> None:
        """A slide with only a full-bleed background picture is empty."""
        prs = _blank_prs()
        slide = _add_blank_slide(prs)
        img_bytes = _tiny_png()
        with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as f:
            f.write(img_bytes)
            tmp_path = f.name
        slide.shapes.add_picture(
            tmp_path, Emu(0), Emu(0),
            width=Emu(_SLIDE_W), height=Emu(_SLIDE_H),
        )
        result = check(prs)
        empty = [i for i in result.issues if i.issue_type == "empty_slide"]
        assert len(empty) == 1

    def test_slide_with_small_image_not_empty(self) -> None:
        """A non-background picture counts as content."""
        prs = _blank_prs()
        slide = _add_blank_slide(prs)
        _add_logo_picture(slide)
        result = check(prs)
        assert not any(i.issue_type == "empty_slide" for i in result.issues)


# ---------------------------------------------------------------------------
# 11. layout_qa — slides_checked metadata
# ---------------------------------------------------------------------------


class TestLayoutQAMetadata:
    def test_slides_checked_matches(self) -> None:
        prs = _blank_prs()
        for _ in range(4):
            slide = _add_blank_slide(prs)
            _add_text_run(slide, "Content")
        result = check(prs)
        assert result.slides_checked == 4

    def test_passed_true_when_clean(self) -> None:
        prs = _blank_prs()
        slide = _add_blank_slide(prs)
        _add_text_run(slide, "Clean content")
        result = check(prs)
        assert result.passed is True
        assert result.issues == []
