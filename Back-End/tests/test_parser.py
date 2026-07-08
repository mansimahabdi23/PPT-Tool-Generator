"""Tests for the expanded deterministic PPTX parser.

Fixture categories
------------------
1. Committed sample deck — broad smoke test across all 6 layouts
2. Text placeholder fixture — role detection, font metadata, heading claims
3. Body text fixture — bullet and statistic claim extraction
4. Table fixture — cell extraction, header_row, table_cell claims
5. Image fixture — image presence, content_type, geometry
6. Geometry fixture — slide-relative fraction correctness

All fixture PPTXs are generated in-memory with python-pptx so tests have
no external file dependencies beyond the committed sample deck.
"""

from __future__ import annotations

import io
import struct
import zlib
from dataclasses import asdict
from pathlib import Path

import pytest
from pptx import Presentation
from pptx.util import Inches, Pt

from app.services.parser import (
    ParsedDeck,
    ParsedSlide,
    parse,
)

SAMPLE_PPTX = Path(__file__).parent.parent / "out" / "samples_imocha_template.pptx"


# ---------------------------------------------------------------------------
# Helpers — build in-memory fixture PPTXs
# ---------------------------------------------------------------------------

def _save(prs: Presentation) -> Path:
    """Save presentation to a temp BytesIO-backed path string for parse()."""
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf  # type: ignore[return-value]  # parse() accepts Path|str|file-like


def _pptx_with_placeholder_text(title: str, body: str) -> io.BytesIO:
    """Slide layout[1] = 'Title and Content' — real title (idx 0) + body (idx 1)."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    slide.placeholders[1].text = body
    return _save(prs)


def _pptx_with_body_stats() -> io.BytesIO:
    """Blank slide with a manual text box containing a statistic string."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # blank
    txb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
    tf = txb.text_frame
    run = tf.paragraphs[0].add_run()
    run.text = "Hiring quality improved 34%, saving $1.2M annually"
    run.font.size = Pt(12)
    return _save(prs)


def _pptx_with_table() -> io.BytesIO:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    tbl = slide.shapes.add_table(3, 2, Inches(1), Inches(1), Inches(6), Inches(3)).table
    tbl.cell(0, 0).text = "Metric"
    tbl.cell(0, 1).text = "Value"
    tbl.cell(1, 0).text = "Hires"
    tbl.cell(1, 1).text = "120"
    tbl.cell(2, 0).text = "Pass rate"
    tbl.cell(2, 1).text = "87%"
    return _save(prs)


def _make_1x1_png() -> bytes:
    """Build a minimal valid 1×1 white RGB PNG without external deps."""
    def chunk(name: bytes, data: bytes) -> bytes:
        length = struct.pack(">I", len(data))
        crc = struct.pack(">I", zlib.crc32(name + data) & 0xFFFF_FFFF)
        return length + name + data + crc

    ihdr = struct.pack(">IIBBBBB", 1, 1, 8, 2, 0, 0, 0)   # 1x1, 8-bit RGB
    idat = zlib.compress(b"\x00\xFF\xFF\xFF")               # filter=none, white pixel
    return (
        b"\x89PNG\r\n\x1a\n"
        + chunk(b"IHDR", ihdr)
        + chunk(b"IDAT", idat)
        + chunk(b"IEND", b"")
    )


def _pptx_with_image() -> io.BytesIO:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    png_bytes = _make_1x1_png()
    slide.shapes.add_picture(
        io.BytesIO(png_bytes),
        Inches(1), Inches(1), Inches(4), Inches(3),
    )
    return _save(prs)


def _parse_buf(buf: io.BytesIO, name: str = "fixture.pptx") -> ParsedDeck:
    """Call parse() on an in-memory BytesIO buffer."""
    buf.seek(0)
    prs_tmp = Presentation(buf)
    tmp = io.BytesIO()
    prs_tmp.save(tmp)
    tmp.seek(0)
    # parse() expects a Path; write to a true temp file
    import tempfile, os
    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
        f.write(tmp.read())
        f.flush()
        path = Path(f.name)
    try:
        deck = parse(path, name)
    finally:
        os.unlink(path)
    return deck


# ---------------------------------------------------------------------------
# 1. Committed sample deck — smoke tests
# ---------------------------------------------------------------------------

def test_sample_deck_exists() -> None:
    assert SAMPLE_PPTX.exists(), f"Sample PPTX missing: {SAMPLE_PPTX}"


def test_sample_deck_slide_count() -> None:
    deck = parse(SAMPLE_PPTX, "sample.pptx")
    assert deck.slide_count == 6
    assert len(deck.slides) == 6


def test_sample_deck_slide_types_at_boundaries() -> None:
    deck = parse(SAMPLE_PPTX, "sample.pptx")
    assert deck.slides[0].slide_type.value == "title"
    assert deck.slides[-1].slide_type.value == "closing"


def test_sample_deck_all_slides_have_text_blocks() -> None:
    deck = parse(SAMPLE_PPTX, "sample.pptx")
    for slide in deck.slides:
        assert len(slide.text_blocks) > 0, f"slide {slide.index} has no text blocks"


def test_sample_deck_all_slides_have_nonempty_title() -> None:
    deck = parse(SAMPLE_PPTX, "sample.pptx")
    for slide in deck.slides:
        assert slide.title, f"slide {slide.index} has empty title"


def test_sample_deck_has_images() -> None:
    """Template slides include logo images."""
    deck = parse(SAMPLE_PPTX, "sample.pptx")
    total_images = sum(len(s.images) for s in deck.slides)
    assert total_images > 0


def test_sample_deck_claims_have_correct_slide_index() -> None:
    deck = parse(SAMPLE_PPTX, "sample.pptx")
    for slide in deck.slides:
        for claim in slide.claims:
            assert claim.source.slide_index == slide.index, (
                f"claim on slide {slide.index} has source.slide_index="
                f"{claim.source.slide_index}"
            )


def test_sample_deck_geometry_fractions_in_range() -> None:
    deck = parse(SAMPLE_PPTX, "sample.pptx")
    for slide in deck.slides:
        for blk in slide.text_blocks:
            g = blk.geometry
            assert 0.0 <= g.left_pct <= 1.05, f"left_pct out of range: {g.left_pct}"
            assert 0.0 <= g.top_pct <= 1.05, f"top_pct out of range: {g.top_pct}"
            assert 0.0 <= g.width_pct <= 1.05, f"width_pct out of range: {g.width_pct}"
            assert 0.0 <= g.height_pct <= 1.05, f"height_pct out of range: {g.height_pct}"


# ---------------------------------------------------------------------------
# 2. Text placeholder fixture — role, font metadata, heading claim
# ---------------------------------------------------------------------------

def test_placeholder_title_role() -> None:
    buf = _pptx_with_placeholder_text("Test Title", "First bullet")
    deck = _parse_buf(buf)
    slide = deck.slides[0]
    title_blocks = [b for b in slide.text_blocks if b.shape_role == "title"]
    assert len(title_blocks) >= 1, "No text block with role='title' found"
    assert title_blocks[0].full_text == "Test Title"


def test_placeholder_title_convenience_field() -> None:
    buf = _pptx_with_placeholder_text("Slide Title", "Some body text")
    deck = _parse_buf(buf)
    assert deck.slides[0].title == "Slide Title"


def test_placeholder_body_in_body_items() -> None:
    buf = _pptx_with_placeholder_text("My Title", "Bullet one\nBullet two")
    deck = _parse_buf(buf)
    body = deck.slides[0].body_items
    assert any("Bullet one" in b for b in body)
    assert any("Bullet two" in b for b in body)


def test_heading_claim_from_title_placeholder() -> None:
    buf = _pptx_with_placeholder_text("Section Header", "Some content")
    deck = _parse_buf(buf)
    slide = deck.slides[0]
    heading_claims = [c for c in slide.claims if c.claim_type == "heading"]
    assert len(heading_claims) >= 1, "Expected at least one heading claim"
    assert heading_claims[0].text == "Section Header"


# ---------------------------------------------------------------------------
# 3. Body text / claim extraction
# ---------------------------------------------------------------------------

def test_statistic_claim_detected() -> None:
    buf = _pptx_with_body_stats()
    deck = _parse_buf(buf)
    slide = deck.slides[0]
    stat_claims = [c for c in slide.claims if c.claim_type == "statistic"]
    assert len(stat_claims) >= 1, "Expected at least one statistic claim"
    assert "34%" in stat_claims[0].text


def test_statistic_source_provenance() -> None:
    buf = _pptx_with_body_stats()
    deck = _parse_buf(buf)
    slide = deck.slides[0]
    for claim in slide.claims:
        assert claim.source.slide_index == 0
        assert isinstance(claim.source.shape_id, int)
        assert claim.source.shape_id > 0


# ---------------------------------------------------------------------------
# 4. Table extraction
# ---------------------------------------------------------------------------

def test_table_row_col_counts() -> None:
    buf = _pptx_with_table()
    deck = _parse_buf(buf)
    slide = deck.slides[0]
    assert len(slide.tables) == 1
    tbl = slide.tables[0]
    assert tbl.row_count == 3
    assert tbl.col_count == 2


def test_table_header_row() -> None:
    buf = _pptx_with_table()
    deck = _parse_buf(buf)
    tbl = deck.slides[0].tables[0]
    assert tbl.header_row == ["Metric", "Value"]


def test_table_cell_values_and_provenance() -> None:
    buf = _pptx_with_table()
    deck = _parse_buf(buf)
    slide = deck.slides[0]
    tbl = slide.tables[0]
    # Find the "87%" cell
    cell_87 = next((c for c in tbl.cells if c.text == "87%"), None)
    assert cell_87 is not None, "87% cell not found"
    assert cell_87.row == 2
    assert cell_87.col == 1
    assert not cell_87.is_header


def test_table_cell_claims_extracted() -> None:
    buf = _pptx_with_table()
    deck = _parse_buf(buf)
    slide = deck.slides[0]
    tc_claims = [c for c in slide.claims if c.claim_type == "table_cell"]
    texts = {c.text for c in tc_claims}
    assert "Metric" in texts
    assert "87%" in texts
    # All table_cell claims carry row/col provenance
    for c in tc_claims:
        assert c.source.row is not None
        assert c.source.col is not None


# ---------------------------------------------------------------------------
# 5. Image extraction
# ---------------------------------------------------------------------------

def test_image_present() -> None:
    buf = _pptx_with_image()
    deck = _parse_buf(buf)
    slide = deck.slides[0]
    assert len(slide.images) >= 1


def test_image_content_type() -> None:
    buf = _pptx_with_image()
    deck = _parse_buf(buf)
    img = deck.slides[0].images[0]
    assert img.content_type.startswith("image/")


def test_image_bytes_size() -> None:
    buf = _pptx_with_image()
    deck = _parse_buf(buf)
    img = deck.slides[0].images[0]
    assert img.bytes_size > 0


def test_image_not_background_when_small() -> None:
    """The 4"×3" image on a 10"×7.5" slide covers < 90% — not a background."""
    buf = _pptx_with_image()
    deck = _parse_buf(buf)
    img = deck.slides[0].images[0]
    assert not img.is_background
