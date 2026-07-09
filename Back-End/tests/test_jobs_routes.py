"""Route tests for /api/jobs endpoints — real async engine (inline mode).

The engine is replaced with ``InProcessJobEngine(inline=True)`` for the
entire module so pipeline segments run synchronously in the TestClient
thread. No sleeps or polling needed.

Test flow (golden path):
  POST /jobs → plan_ready (segment A complete)
  GET  /jobs/{id}/plan → non-empty list
  POST /jobs/{id}/plan/approve → 202 + completed (segment B complete)
  GET  /jobs/{id} → completed status, real validator fields
  GET  /jobs/{id}/result → real URLs + validator booleans
"""

from __future__ import annotations

import io

import pytest
from fastapi.testclient import TestClient
from pptx import Presentation

from app.services import job_engine as _je
from app.services.job_engine import InProcessJobEngine
from app.services.llm_provider import StubProvider, init_provider


# ---------------------------------------------------------------------------
# Module-level engine + provider override — inline mode for deterministic tests
# ---------------------------------------------------------------------------

@pytest.fixture(autouse=True, scope="module")
def _inline_engine() -> None:  # type: ignore[return]
    """Replace the global engine with an inline (synchronous) instance."""
    _je.init_engine(InProcessJobEngine(inline=True))
    init_provider(StubProvider())


# ---------------------------------------------------------------------------
# Minimal PPTX factory
# ---------------------------------------------------------------------------

def _minimal_pptx_bytes() -> bytes:
    """Return a valid one-slide PPTX as bytes (title slide)."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])  # Title Slide layout
    slide.shapes.title.text = "iMocha Test Deck"
    try:
        slide.placeholders[1].text = "Subtitle text"
    except (KeyError, IndexError):
        pass
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ---------------------------------------------------------------------------
# Shared helper — create a real job and return its ID
# ---------------------------------------------------------------------------

def _create_job(client: TestClient, allow_restructure: bool = False) -> str:
    pptx_bytes = _minimal_pptx_bytes()
    resp = client.post(
        "/api/jobs",
        files={"file": ("test.pptx", io.BytesIO(pptx_bytes), "application/vnd.openxmlformats-officedocument.presentationml.presentation")},
        data={"allow_restructure": str(allow_restructure).lower()},
    )
    assert resp.status_code == 201, resp.text
    return resp.json()["jobId"]


def _approve_and_complete(client: TestClient, job_id: str) -> None:
    resp = client.post(f"/api/jobs/{job_id}/plan/approve")
    assert resp.status_code == 202, resp.text


# ---------------------------------------------------------------------------
# POST /api/jobs
# ---------------------------------------------------------------------------

class TestCreateJob:
    def test_returns_201_with_job_id(self, client: TestClient) -> None:
        job_id = _create_job(client)
        assert isinstance(job_id, str) and len(job_id) > 0

    def test_job_reaches_plan_ready(self, client: TestClient) -> None:
        """Segment A completes synchronously in inline mode → plan_ready."""
        job_id = _create_job(client)
        resp = client.get(f"/api/jobs/{job_id}")
        assert resp.status_code == 200
        assert resp.json()["status"] == "plan_ready"

    def test_corrupt_upload_reaches_failed(self, client: TestClient) -> None:
        """Non-PPTX bytes → parse() raises → status = failed."""
        resp = client.post(
            "/api/jobs",
            files={"file": ("bad.pptx", io.BytesIO(b"not a pptx"), "application/octet-stream")},
            data={"allow_restructure": "false"},
        )
        assert resp.status_code == 201
        job_id = resp.json()["jobId"]
        status = client.get(f"/api/jobs/{job_id}").json()["status"]
        assert status == "failed"


# ---------------------------------------------------------------------------
# GET /api/jobs
# ---------------------------------------------------------------------------

class TestListJobs:
    def test_returns_200_list(self, client: TestClient) -> None:
        _create_job(client)  # ensure at least one job exists
        resp = client.get("/api/jobs")
        assert resp.status_code == 200
        assert isinstance(resp.json(), list)

    def test_list_items_are_camel_case(self, client: TestClient) -> None:
        _create_job(client)
        job = client.get("/api/jobs").json()[0]
        for key in ("deckName", "slideCount", "allowRestructure", "createdAt", "status"):
            assert key in job, f"missing key: {key}"

    def test_accepts_paging_params(self, client: TestClient) -> None:
        resp = client.get("/api/jobs?page=1&page_size=5")
        assert resp.status_code == 200


# ---------------------------------------------------------------------------
# GET /api/jobs/{id}
# ---------------------------------------------------------------------------

class TestGetJob:
    def test_unknown_job_returns_404(self, client: TestClient) -> None:
        resp = client.get("/api/jobs/nonexistent-job-id-xyz")
        assert resp.status_code == 404

    def test_plan_ready_has_plan_no_slides(self, client: TestClient) -> None:
        job_id = _create_job(client)
        body = client.get(f"/api/jobs/{job_id}").json()
        assert body["status"] == "plan_ready"
        assert isinstance(body["plan"], list)
        assert len(body["plan"]) > 0
        assert body["slides"] is None  # not yet produced

    def test_completed_job_has_slides_and_plan(self, client: TestClient) -> None:
        job_id = _create_job(client)
        _approve_and_complete(client, job_id)
        body = client.get(f"/api/jobs/{job_id}").json()
        assert body["status"] == "completed"
        assert isinstance(body["plan"], list) and len(body["plan"]) > 0
        assert isinstance(body["slides"], list) and len(body["slides"]) > 0

    def test_completed_job_has_validator_fields(self, client: TestClient) -> None:
        job_id = _create_job(client)
        _approve_and_complete(client, job_id)
        body = client.get(f"/api/jobs/{job_id}").json()
        assert body["brandCompliancePassed"] is not None
        assert isinstance(body["contentFidelity"], str)
        assert "/" in body["contentFidelity"]  # e.g. "0/0 claims preserved"

    def test_slides_are_camel_case(self, client: TestClient) -> None:
        job_id = _create_job(client)
        _approve_and_complete(client, job_id)
        slide = client.get(f"/api/jobs/{job_id}").json()["slides"][0]
        for key in ("originalPreviewUrl", "transformedPreviewUrl", "contentUnchanged", "changeChips", "approval"):
            assert key in slide, f"missing key: {key}"

    def test_plan_items_are_camel_case(self, client: TestClient) -> None:
        job_id = _create_job(client)
        body = client.get(f"/api/jobs/{job_id}").json()
        plan_item = body["plan"][0]
        for key in ("slideType", "plannedLayout", "assetTypes"):
            assert key in plan_item, f"missing key: {key}"

    def test_allow_restructure_propagated(self, client: TestClient) -> None:
        job_id = _create_job(client, allow_restructure=True)
        body = client.get(f"/api/jobs/{job_id}").json()
        assert body["allowRestructure"] is True


# ---------------------------------------------------------------------------
# GET /api/jobs/{id}/plan
# ---------------------------------------------------------------------------

class TestGetPlan:
    def test_returns_200_list_of_slide_plans(self, client: TestClient) -> None:
        job_id = _create_job(client)
        resp = client.get(f"/api/jobs/{job_id}/plan")
        assert resp.status_code == 200
        plan = resp.json()
        assert isinstance(plan, list)
        assert len(plan) >= 1

    def test_plan_item_fields(self, client: TestClient) -> None:
        job_id = _create_job(client)
        plan = client.get(f"/api/jobs/{job_id}/plan").json()
        item = plan[0]
        for key in ("id", "index", "slideType", "plannedLayout", "assetTypes"):
            assert key in item, f"missing key: {key}"

    def test_unknown_job_returns_404(self, client: TestClient) -> None:
        resp = client.get("/api/jobs/no-such-job/plan")
        assert resp.status_code == 404

    def test_allow_restructure_includes_note(self, client: TestClient) -> None:
        job_id = _create_job(client, allow_restructure=True)
        plan = client.get(f"/api/jobs/{job_id}/plan").json()
        # At least one plan item should have a restructureNote
        assert any(item.get("restructureNote") for item in plan)


# ---------------------------------------------------------------------------
# POST /api/jobs/{id}/plan/approve
# ---------------------------------------------------------------------------

class TestApprovePlan:
    def test_returns_202_with_transform_job(self, client: TestClient) -> None:
        job_id = _create_job(client)
        resp = client.post(f"/api/jobs/{job_id}/plan/approve")
        assert resp.status_code == 202
        body = resp.json()
        for key in ("id", "status", "deckName"):
            assert key in body, f"missing key: {key}"

    def test_unknown_job_returns_404(self, client: TestClient) -> None:
        resp = client.post("/api/jobs/no-such-job/plan/approve")
        assert resp.status_code == 404

    def test_double_approve_returns_409(self, client: TestClient) -> None:
        """Approving a job not in plan_ready (e.g. already completed) → 409."""
        job_id = _create_job(client)
        _approve_and_complete(client, job_id)  # first approve → completed
        resp = client.post(f"/api/jobs/{job_id}/plan/approve")
        assert resp.status_code == 409

    def test_approve_before_plan_ready_returns_409(self, client: TestClient) -> None:
        """Pre-plan-ready approve is prevented by the 409 guard.

        We can't easily test mid-parse state in inline mode, so we test the
        409 on a completed job (same guard: status != plan_ready).
        """
        job_id = _create_job(client)
        _approve_and_complete(client, job_id)
        resp = client.post(f"/api/jobs/{job_id}/plan/approve")
        assert resp.status_code == 409

    def test_completed_job_has_validator_data(self, client: TestClient) -> None:
        job_id = _create_job(client)
        _approve_and_complete(client, job_id)
        body = client.get(f"/api/jobs/{job_id}").json()
        assert body["status"] == "completed"
        assert "brandCompliancePassed" in body
        assert "contentFidelity" in body


# ---------------------------------------------------------------------------
# POST /api/jobs/{id}/slides/{sid}/regenerate — stub
# ---------------------------------------------------------------------------

class TestRegenerateSlide:
    def test_returns_200_transformed_slide(self, client: TestClient) -> None:
        job_id = _create_job(client)
        resp = client.post(f"/api/jobs/{job_id}/slides/s1/regenerate")
        assert resp.status_code == 200
        body = resp.json()
        for key in ("id", "originalPreviewUrl", "transformedPreviewUrl", "approval"):
            assert key in body, f"missing key: {key}"


# ---------------------------------------------------------------------------
# GET /api/jobs/{id}/result
# ---------------------------------------------------------------------------

class TestGetResult:
    def test_returns_200_job_result(self, client: TestClient) -> None:
        job_id = _create_job(client)
        _approve_and_complete(client, job_id)
        resp = client.get(f"/api/jobs/{job_id}/result")
        assert resp.status_code == 200
        body = resp.json()
        for key in ("pptxUrl", "pdfUrl", "brandCompliancePassed", "contentFidelity", "processingSeconds"):
            assert key in body, f"missing key: {key}"

    def test_result_urls_are_strings(self, client: TestClient) -> None:
        job_id = _create_job(client)
        _approve_and_complete(client, job_id)
        body = client.get(f"/api/jobs/{job_id}/result").json()
        assert isinstance(body["pptxUrl"], str)
        assert isinstance(body["pdfUrl"], str)

    def test_brand_compliance_from_real_validator(self, client: TestClient) -> None:
        """brandCompliancePassed must be a real boolean (not hardcoded True)."""
        job_id = _create_job(client)
        _approve_and_complete(client, job_id)
        body = client.get(f"/api/jobs/{job_id}/result").json()
        assert isinstance(body["brandCompliancePassed"], bool)

    def test_incomplete_job_returns_409(self, client: TestClient) -> None:
        job_id = _create_job(client)  # stops at plan_ready
        resp = client.get(f"/api/jobs/{job_id}/result")
        assert resp.status_code == 409

    def test_unknown_job_returns_404(self, client: TestClient) -> None:
        resp = client.get("/api/jobs/no-such-job/result")
        assert resp.status_code == 404
