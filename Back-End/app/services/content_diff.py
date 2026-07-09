"""Deterministic content-diff validator — no LLM.

Verifies that every atomic claim (factual statement, statistic, table cell) from
the source ParsedDeck is present, unaltered, in the output Presentation.

"Present" means the normalized claim text is a substring of the normalized output
text. Normalization collapses whitespace and lower-cases — it does NOT change words,
numbers, or special characters, so statistics like "34%" must survive verbatim.

Public API
----------
diff(source: ParsedDeck, output_prs: Presentation) -> ContentDiffResult
"""

from __future__ import annotations

import re
from dataclasses import dataclass, field

from pptx import Presentation

from app.services.parser import AtomicClaim, ParsedDeck, SourceRef


# ---------------------------------------------------------------------------
# Result types
# ---------------------------------------------------------------------------


@dataclass
class MissingClaim:
    """A source claim not found in the output."""

    text: str
    claim_type: str  # "heading" | "statistic" | "bullet" | "table_cell"
    source_slide: int
    source_shape: str


@dataclass
class ContentDiffResult:
    passed: bool
    total_claims: int
    preserved_claims: int
    missing_claims: list[MissingClaim] = field(default_factory=list)
    fidelity_str: str = ""  # e.g. "43/45 claims preserved"

    def __post_init__(self) -> None:
        if not self.fidelity_str:
            self.fidelity_str = (
                f"{self.preserved_claims}/{self.total_claims} claims preserved"
            )


# ---------------------------------------------------------------------------
# Text extraction from output Presentation
# ---------------------------------------------------------------------------


def _extract_output_text(prs: Presentation) -> str:
    """Collect all text from the output deck into one normalized blob."""
    parts: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                try:
                    parts.append(shape.text_frame.text)
                except Exception:
                    pass
            if getattr(shape, "has_table", False):
                try:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            parts.append(cell.text_frame.text)
                except Exception:
                    pass
    return " ".join(parts)


def _normalize(text: str) -> str:
    """Collapse whitespace and lower-case; preserve all other characters."""
    return re.sub(r"\s+", " ", text.strip()).lower()


# ---------------------------------------------------------------------------
# Public API
# ---------------------------------------------------------------------------


def diff(source: ParsedDeck, output_prs: Presentation) -> ContentDiffResult:
    """Check every source claim against the output Presentation text.

    Parameters
    ----------
    source     : ParsedDeck produced by parser.parse() on the uploaded file
    output_prs : composed output python-pptx Presentation

    Returns
    -------
    ContentDiffResult — passed=True iff all claims are preserved.
    """
    all_claims: list[AtomicClaim] = []
    for slide in source.slides:
        all_claims.extend(slide.claims)

    total = len(all_claims)
    if total == 0:
        return ContentDiffResult(
            passed=True,
            total_claims=0,
            preserved_claims=0,
            fidelity_str="0/0 claims preserved (source has no claims)",
        )

    output_blob = _normalize(_extract_output_text(output_prs))

    missing: list[MissingClaim] = []
    for claim in all_claims:
        needle = _normalize(claim.text)
        if needle and needle not in output_blob:
            missing.append(MissingClaim(
                text=claim.text,
                claim_type=claim.claim_type,
                source_slide=claim.source.slide_index,
                source_shape=claim.source.shape_name,
            ))

    preserved = total - len(missing)
    return ContentDiffResult(
        passed=len(missing) == 0,
        total_claims=total,
        preserved_claims=preserved,
        missing_claims=missing,
        fidelity_str=f"{preserved}/{total} claims preserved",
    )
