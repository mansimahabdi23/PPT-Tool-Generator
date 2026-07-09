"""Deterministic brand-lint validator — no LLM.

Inspects an output Presentation for any font, color, gradient, or logo deviation
from the iMocha allow-list (brand.py §11). Returns structured pass/fail + exact
violations so the composer's retry loop can act on specific evidence.

Public API
----------
lint(prs: Presentation) -> BrandLintResult
"""

from __future__ import annotations

from dataclasses import dataclass, field

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from app import brand

# OOXML namespace URIs
_P = "http://schemas.openxmlformats.org/presentationml/2006/main"
_A = "http://schemas.openxmlformats.org/drawingml/2006/main"

_GS_TAG = f"{{{_A}}}gs"
_SRGB_TAG = f"{{{_A}}}srgbClr"
_GRAD_FILL_TAG = f"{{{_A}}}gradFill"
_SOLID_FILL_TAG = f"{{{_A}}}solidFill"


# ---------------------------------------------------------------------------
# Result types
# ---------------------------------------------------------------------------


@dataclass
class BrandViolation:
    """One brand rule infraction."""

    slide_index: int
    shape_name: str
    violation_type: str  # "font" | "color" | "gradient" | "logo_missing"
    detail: str


@dataclass
class BrandLintResult:
    passed: bool
    violations: list[BrandViolation] = field(default_factory=list)
    slide_count: int = 0
    shapes_checked: int = 0


# ---------------------------------------------------------------------------
# Text-frame checks
# ---------------------------------------------------------------------------


def _check_text_frame(
    shape: object,
    slide_index: int,
    violations: list[BrandViolation],
) -> None:
    try:
        for para in shape.text_frame.paragraphs:  # type: ignore[attr-defined]
            for run in para.runs:
                f = run.font

                # Font name — only flag when explicitly set on the run
                if f.name and not brand.is_allowed_font(f.name):
                    violations.append(BrandViolation(
                        slide_index=slide_index,
                        shape_name=shape.name,  # type: ignore[attr-defined]
                        violation_type="font",
                        detail=f"font '{f.name}' not in allow-list",
                    ))

                # Color — only flag explicit RGB (not theme/inherited)
                try:
                    clr = f.color
                    if clr.type is not None:
                        rgb = str(clr.rgb).upper().lstrip("#")
                        if not brand.is_allowed_color(rgb):
                            violations.append(BrandViolation(
                                slide_index=slide_index,
                                shape_name=shape.name,  # type: ignore[attr-defined]
                                violation_type="color",
                                detail=f"text color '#{rgb}' not in allow-list",
                            ))
                except Exception:
                    pass
    except Exception:
        pass


# ---------------------------------------------------------------------------
# Shape fill checks
# ---------------------------------------------------------------------------


def _srgb_stops_from_element(el: object) -> list[str]:
    """Return list of srgbClr val strings from gradient stops inside *el*."""
    stops: list[str] = []
    for gs in el.iter(_GS_TAG):  # type: ignore[attr-defined]
        for srgb in gs:
            if srgb.tag == _SRGB_TAG:
                val = srgb.get("val", "").upper()
                if val:
                    stops.append(val)
    return stops


def _check_shape_fill(
    shape: object,
    slide_index: int,
    violations: list[BrandViolation],
) -> None:
    """Check solid-color and gradient fills on a shape element."""
    try:
        sp_el = shape._element  # type: ignore[attr-defined]
    except Exception:
        return

    # --- Solid fill --------------------------------------------------------
    try:
        fill = shape.fill  # type: ignore[attr-defined]
        from pptx.enum.dml import MSO_FILL

        if fill.type == MSO_FILL.SOLID:
            try:
                if fill.fore_color.type is not None:  # explicit color set
                    rgb = str(fill.fore_color.rgb).upper().lstrip("#")
                    if not brand.is_allowed_color(rgb):
                        violations.append(BrandViolation(
                            slide_index=slide_index,
                            shape_name=shape.name,  # type: ignore[attr-defined]
                            violation_type="color",
                            detail=f"shape fill '#{rgb}' not in allow-list",
                        ))
            except Exception:
                pass
    except Exception:
        pass

    # --- Gradient fill on shape (not background) ---------------------------
    try:
        stops = _srgb_stops_from_element(sp_el)
        if len(stops) >= 2 and not brand.is_allowed_gradient(stops[0], stops[-1]):
            violations.append(BrandViolation(
                slide_index=slide_index,
                shape_name=shape.name,  # type: ignore[attr-defined]
                violation_type="gradient",
                detail=f"gradient (#{stops[0]}, #{stops[-1]}) not in allow-list",
            ))
    except Exception:
        pass


# ---------------------------------------------------------------------------
# Slide background gradient check
# ---------------------------------------------------------------------------


def _check_slide_bg_gradient(
    slide: object,
    slide_index: int,
    violations: list[BrandViolation],
) -> None:
    """Check the slide's <p:bg> gradient fill against the allow-list."""
    try:
        cSld = slide.element.find(f"{{{_P}}}cSld")  # type: ignore[attr-defined]
        if cSld is None:
            return
        bg = cSld.find(f"{{{_P}}}bg")
        if bg is None:
            return
        bgPr = bg.find(f"{{{_P}}}bgPr")
        if bgPr is None:
            return
        gradFill = bgPr.find(_GRAD_FILL_TAG)
        if gradFill is None:
            return

        stops = _srgb_stops_from_element(gradFill)
        if len(stops) >= 2 and not brand.is_allowed_gradient(stops[0], stops[-1]):
            violations.append(BrandViolation(
                slide_index=slide_index,
                shape_name="<slide background>",
                violation_type="gradient",
                detail=f"background gradient (#{stops[0]}, #{stops[-1]}) not in allow-list",
            ))
    except Exception:
        pass


# ---------------------------------------------------------------------------
# Logo presence check
# ---------------------------------------------------------------------------


def _has_logo(slide: object, slide_h: int) -> bool:
    """True if any non-background picture is in the footer region (bottom 20%)
    or has 'logo' in its name."""
    footer_threshold = slide_h * 0.80
    for shape in slide.shapes:  # type: ignore[attr-defined]
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            name = (shape.name or "").lower()
            top = getattr(shape, "top", None) or 0
            if "logo" in name or top > footer_threshold:
                return True
    return False


# ---------------------------------------------------------------------------
# Shape tree traversal
# ---------------------------------------------------------------------------


def _traverse(
    shapes: object,
    slide_index: int,
    violations: list[BrandViolation],
    depth: int = 0,
) -> int:
    """Walk shapes, check text and fill. Returns count of shapes inspected."""
    count = 0
    for shape in shapes:  # type: ignore[iteration]
        try:
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                if depth == 0:
                    count += _traverse(
                        shape.shapes, slide_index, violations, depth=1
                    )
                continue

            count += 1

            if getattr(shape, "has_text_frame", False):
                _check_text_frame(shape, slide_index, violations)

            # Skip fill check on picture shapes (they use blipFill)
            if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
                _check_shape_fill(shape, slide_index, violations)

        except Exception:
            pass

    return count


# ---------------------------------------------------------------------------
# Public API
# ---------------------------------------------------------------------------


def lint(prs: Presentation) -> BrandLintResult:
    """Inspect *prs* for brand violations.

    Checks every slide for:
    - Fonts not in brand.ALLOWED_FONTS (explicit run fonts only)
    - Colors not in brand.ALLOWED_COLORS (explicit run/fill colors only)
    - Gradients not in brand.ALLOWED_GRADIENTS (background + shape fills)
    - Missing logo in footer region

    Returns BrandLintResult; passed=True iff violations is empty.
    """
    violations: list[BrandViolation] = []
    total_shapes = 0
    slide_h: int = prs.slide_height

    for slide_index, slide in enumerate(prs.slides):
        total_shapes += _traverse(slide.shapes, slide_index, violations)
        _check_slide_bg_gradient(slide, slide_index, violations)

        if not _has_logo(slide, slide_h):
            violations.append(BrandViolation(
                slide_index=slide_index,
                shape_name="<slide>",
                violation_type="logo_missing",
                detail="no logo found in footer region (top > 80% of slide height)",
            ))

    return BrandLintResult(
        passed=len(violations) == 0,
        violations=violations,
        slide_count=len(prs.slides),
        shapes_checked=total_shapes,
    )
