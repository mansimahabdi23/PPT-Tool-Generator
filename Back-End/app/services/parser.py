"""Deterministic PPTX parser — no AI.

Normalizes an uploaded deck into a structured slide model carrying:
  · text blocks with full font metadata (name, size, bold, italic, color hex)
  · tables with per-cell provenance (slide, shape, row, col)
  · images with geometry and content-type
  · atomic claims (factual statements / statistics) with source-ref for fidelity tracking

The downstream AI Analyze & Plan agent reads this model to understand content and
classify slides. It never re-parses the raw PPTX — provenance is attached here so
the content-diff validator can verify every claim survived transformation.

Public API
----------
parse(pptx_path, deck_name) -> ParsedDeck
"""

from __future__ import annotations

import re
from dataclasses import dataclass, field
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN

from app.models.enums import SlideType

# ---------------------------------------------------------------------------
# Claim-detection pattern — deterministic regex, no AI
# ---------------------------------------------------------------------------

_STAT_RE = re.compile(
    r"""
    (?:
        (?<!\w)\d[\d,]*\.?\d*\s*%              # 34%, 12.5%
      | (?<!\w)\d[\d,]*\.?\d*\s*[xX×]\b        # 2x, 3X, 2×
      | [$£€¥]\s*\d[\d,]*\.?\d*                 # $1.2M, £50
      | (?<!\w)\d[\d,]*\.?\d*\s*[kKmMbBtT]\b   # 5M, 2.1B, 120K
      | (?<!\w)\d+\s*/\s*\d+                    # 3/5
    )
    """,
    re.VERBOSE,
)

# ---------------------------------------------------------------------------
# Data models
# ---------------------------------------------------------------------------


@dataclass
class ShapeGeometry:
    """Bounding box in EMU and slide-relative fractions (0.0–1.0)."""

    left_emu: int
    top_emu: int
    width_emu: int
    height_emu: int
    left_pct: float
    top_pct: float
    width_pct: float
    height_pct: float


@dataclass
class SourceRef:
    """Exact provenance pointer for one piece of content.

    Used by the content-diff validator to confirm every claim survived
    the transformation step unchanged.
    """

    slide_index: int         # 0-based
    shape_id: int
    shape_name: str
    paragraph_index: int | None = None
    row: int | None = None   # table row
    col: int | None = None   # table column


@dataclass
class TextRun:
    """One font-decorated span within a paragraph."""

    text: str
    bold: bool
    italic: bool
    font_name: str | None
    font_size_pt: float | None
    color_hex: str | None    # "RRGGBB" (explicit RGB only); None = theme/inherited


@dataclass
class TextParagraph:
    index: int
    text: str                # full paragraph text
    runs: list[TextRun]
    level: int               # indent level; 0 = top
    alignment: str           # "left" | "center" | "right" | "justify" | "unknown"


@dataclass
class ParsedTextBlock:
    shape_id: int
    shape_name: str
    geometry: ShapeGeometry
    shape_role: str          # "title" | "body" | "other"
    paragraphs: list[TextParagraph]
    full_text: str           # "\n".join of non-empty paragraph texts


@dataclass
class TableCell:
    row: int
    col: int
    text: str
    is_header: bool          # True for row 0


@dataclass
class ParsedTable:
    shape_id: int
    shape_name: str
    geometry: ShapeGeometry
    row_count: int
    col_count: int
    cells: list[TableCell]
    header_row: list[str]    # texts of row 0 (convenience)


@dataclass
class ParsedImage:
    shape_id: int
    shape_name: str
    geometry: ShapeGeometry
    alt_text: str
    content_type: str        # "image/png", "image/jpeg", etc.
    bytes_size: int
    is_background: bool      # True when ≥90% of slide dimensions


@dataclass
class AtomicClaim:
    """One factual statement / data point with exact source provenance.

    Deterministically extracted — no interpretation, no paraphrasing.
    claim_type is a coarse hint; the AI Analyze & Plan agent assigns
    semantic meaning. The content-diff validator uses source to verify
    every claim survives transformation.
    """

    text: str
    claim_type: str          # "heading" | "statistic" | "bullet" | "table_cell"
    source: SourceRef


@dataclass
class ParsedSlide:
    index: int
    slide_type: SlideType    # heuristic; AI Analyze & Plan agent overrides later
    text_blocks: list[ParsedTextBlock]
    tables: list[ParsedTable]
    images: list[ParsedImage]
    claims: list[AtomicClaim]
    # Convenience fields — kept for composer backward compatibility
    title: str
    body_items: list[str]


@dataclass
class ParsedDeck:
    name: str
    slide_count: int
    slides: list[ParsedSlide] = field(default_factory=list)


# ---------------------------------------------------------------------------
# Geometry helpers
# ---------------------------------------------------------------------------

def _geom(shape: object, slide_w: int, slide_h: int) -> ShapeGeometry:
    l: int = getattr(shape, "left", None) or 0
    t: int = getattr(shape, "top", None) or 0
    w: int = getattr(shape, "width", None) or 0
    h: int = getattr(shape, "height", None) or 0
    return ShapeGeometry(
        left_emu=l, top_emu=t, width_emu=w, height_emu=h,
        left_pct=l / slide_w if slide_w else 0.0,
        top_pct=t / slide_h if slide_h else 0.0,
        width_pct=w / slide_w if slide_w else 0.0,
        height_pct=h / slide_h if slide_h else 0.0,
    )


# ---------------------------------------------------------------------------
# Text extraction helpers
# ---------------------------------------------------------------------------

_ALIGN_MAP = {
    PP_ALIGN.LEFT: "left",
    PP_ALIGN.CENTER: "center",
    PP_ALIGN.RIGHT: "right",
    PP_ALIGN.JUSTIFY: "justify",
    PP_ALIGN.DISTRIBUTE: "justify",
    PP_ALIGN.THAI_DISTRIBUTE: "justify",
}


def _color_hex(font: object) -> str | None:
    try:
        color = getattr(font, "color", None)
        if color and color.type is not None:
            return str(color.rgb)   # "RRGGBB"
    except Exception:
        pass
    return None


def _parse_runs(para: object) -> list[TextRun]:
    runs: list[TextRun] = []
    for run in para.runs:  # type: ignore[attr-defined]
        f = run.font
        size_pt: float | None = None
        try:
            if f.size:
                size_pt = f.size.pt
        except Exception:
            pass
        runs.append(TextRun(
            text=run.text,
            bold=bool(f.bold),
            italic=bool(f.italic),
            font_name=f.name,
            font_size_pt=size_pt,
            color_hex=_color_hex(f),
        ))
    return runs


def _shape_role(shape: object) -> str:
    if getattr(shape, "is_placeholder", False):
        try:
            idx = shape.placeholder_format.idx  # type: ignore[attr-defined]
            return "title" if idx == 0 else "body"
        except (ValueError, AttributeError):
            pass
    return "other"


def _build_text_block(shape: object, geom: ShapeGeometry) -> ParsedTextBlock:
    role = _shape_role(shape)
    paragraphs: list[TextParagraph] = []
    lines: list[str] = []
    for i, para in enumerate(shape.text_frame.paragraphs):  # type: ignore[attr-defined]
        text = para.text.strip()
        align = _ALIGN_MAP.get(para.alignment, "unknown")
        paragraphs.append(TextParagraph(
            index=i,
            text=text,
            runs=_parse_runs(para),
            level=para.level,
            alignment=align,
        ))
        if text:
            lines.append(text)
    return ParsedTextBlock(
        shape_id=shape.shape_id,  # type: ignore[attr-defined]
        shape_name=shape.name,  # type: ignore[attr-defined]
        geometry=geom,
        shape_role=role,
        paragraphs=paragraphs,
        full_text="\n".join(lines),
    )


# ---------------------------------------------------------------------------
# Table extraction helpers
# ---------------------------------------------------------------------------

def _build_table(shape: object, geom: ShapeGeometry) -> ParsedTable:
    tbl = shape.table  # type: ignore[attr-defined]
    rows, cols = len(tbl.rows), len(tbl.columns)
    cells: list[TableCell] = []
    for r in range(rows):
        for c in range(cols):
            try:
                text = tbl.cell(r, c).text_frame.text.strip()
            except Exception:
                text = ""
            cells.append(TableCell(row=r, col=c, text=text, is_header=(r == 0)))
    header_row = [cell.text for cell in cells if cell.is_header]
    return ParsedTable(
        shape_id=shape.shape_id,  # type: ignore[attr-defined]
        shape_name=shape.name,  # type: ignore[attr-defined]
        geometry=geom,
        row_count=rows,
        col_count=cols,
        cells=cells,
        header_row=header_row,
    )


# ---------------------------------------------------------------------------
# Image extraction helpers
# ---------------------------------------------------------------------------

def _build_image(
    shape: object,
    geom: ShapeGeometry,
    slide_w: int,
    slide_h: int,
) -> ParsedImage:
    try:
        img = shape.image  # type: ignore[attr-defined]
        content_type = img.content_type
        bytes_size = len(img.blob)
    except Exception:
        content_type = "image/unknown"
        bytes_size = 0
    is_bg = (
        geom.width_pct >= 0.90
        and geom.height_pct >= 0.90
        and geom.left_pct <= 0.05
        and geom.top_pct <= 0.05
    )
    return ParsedImage(
        shape_id=shape.shape_id,  # type: ignore[attr-defined]
        shape_name=shape.name,  # type: ignore[attr-defined]
        geometry=geom,
        alt_text=shape.name,  # type: ignore[attr-defined]  # python-pptx doesn't expose alt text natively
        content_type=content_type,
        bytes_size=bytes_size,
        is_background=is_bg,
    )


# ---------------------------------------------------------------------------
# Shape tree traversal (one level of group recursion)
# ---------------------------------------------------------------------------

def _traverse(
    shapes: object,
    slide_index: int,
    slide_w: int,
    slide_h: int,
    text_blocks: list[ParsedTextBlock],
    tables: list[ParsedTable],
    images: list[ParsedImage],
    depth: int = 0,
) -> None:
    for shape in shapes:  # type: ignore[iteration]
        try:
            geom = _geom(shape, slide_w, slide_h)
        except Exception:
            continue

        # Recurse into groups one level deep only
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            if depth == 0:
                try:
                    _traverse(
                        shape.shapes, slide_index, slide_w, slide_h,
                        text_blocks, tables, images, depth=1,
                    )
                except Exception:
                    pass
            continue

        if getattr(shape, "has_text_frame", False):
            try:
                text_blocks.append(_build_text_block(shape, geom))
            except Exception:
                pass
        elif getattr(shape, "has_table", False):
            try:
                tables.append(_build_table(shape, geom))
            except Exception:
                pass
        elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            try:
                images.append(_build_image(shape, geom, slide_w, slide_h))
            except Exception:
                pass


# ---------------------------------------------------------------------------
# Claim extraction — deterministic
# ---------------------------------------------------------------------------

def _claims_from_blocks(
    slide_index: int,
    blocks: list[ParsedTextBlock],
) -> list[AtomicClaim]:
    claims: list[AtomicClaim] = []
    for block in blocks:
        for para in block.paragraphs:
            text = para.text.strip()
            if not text:
                continue
            if block.shape_role == "title":
                ctype = "heading"
            elif _STAT_RE.search(text):
                ctype = "statistic"
            else:
                ctype = "bullet"
            claims.append(AtomicClaim(
                text=text,
                claim_type=ctype,
                source=SourceRef(
                    slide_index=slide_index,
                    shape_id=block.shape_id,
                    shape_name=block.shape_name,
                    paragraph_index=para.index,
                ),
            ))
    return claims


def _claims_from_tables(
    slide_index: int,
    tables: list[ParsedTable],
) -> list[AtomicClaim]:
    claims: list[AtomicClaim] = []
    for tbl in tables:
        for cell in tbl.cells:
            text = cell.text.strip()
            if not text:
                continue
            claims.append(AtomicClaim(
                text=text,
                claim_type="table_cell",
                source=SourceRef(
                    slide_index=slide_index,
                    shape_id=tbl.shape_id,
                    shape_name=tbl.shape_name,
                    row=cell.row,
                    col=cell.col,
                ),
            ))
    return claims


# ---------------------------------------------------------------------------
# Title / body convenience fields (composer backward compat)
# ---------------------------------------------------------------------------

def _derive_title(blocks: list[ParsedTextBlock]) -> str:
    # Priority 1: placeholder title (role="title")
    for b in blocks:
        if b.shape_role == "title" and b.full_text:
            return b.full_text
    # Fallback: text block with largest font size on any run
    best, best_pt = "", 0.0
    for b in blocks:
        for para in b.paragraphs:
            for run in para.runs:
                if run.font_size_pt and run.font_size_pt > best_pt and para.text.strip():
                    best_pt = run.font_size_pt
                    best = b.full_text
    return best


_FOOTER_TOP_PCT = 0.88  # text blocks starting below this are footers/watermarks


def _derive_body(blocks: list[ParsedTextBlock], title_text: str) -> list[str]:
    lines: list[str] = []
    for b in blocks:
        if b.shape_role == "title" or b.full_text == title_text:
            continue
        if b.geometry.top_pct > _FOOTER_TOP_PCT:  # skip footer / page-number zone
            continue
        for para in b.paragraphs:
            line = para.text.strip()
            if line:
                lines.append(line)
    return lines


# ---------------------------------------------------------------------------
# Slide-type heuristic (AI Analyze & Plan agent overrides in Step 5)
# ---------------------------------------------------------------------------

def _assign_type(
    index: int,
    last_index: int,
    title: str,
    body_items: list[str],
) -> SlideType:
    if index == 0:
        return SlideType.title
    if index == last_index:
        return SlideType.closing
    if "agenda" in title.lower():
        return SlideType.agenda
    if not body_items:
        return SlideType.divider
    return SlideType.content


# ---------------------------------------------------------------------------
# Public API
# ---------------------------------------------------------------------------

def parse(pptx_path: Path, deck_name: str) -> ParsedDeck:
    """Open *pptx_path* and return a fully normalized ParsedDeck.

    Parameters
    ----------
    pptx_path : path to the source .pptx file
    deck_name : display name for the deck (typically the original filename)
    """
    prs = Presentation(str(pptx_path))
    slide_w: int = prs.slide_width
    slide_h: int = prs.slide_height
    raw_slides = prs.slides
    last = len(raw_slides) - 1

    parsed_slides: list[ParsedSlide] = []
    for i, slide in enumerate(raw_slides):
        text_blocks: list[ParsedTextBlock] = []
        tables: list[ParsedTable] = []
        images: list[ParsedImage] = []

        _traverse(slide.shapes, i, slide_w, slide_h, text_blocks, tables, images)

        title = _derive_title(text_blocks)
        body = _derive_body(text_blocks, title)
        slide_type = _assign_type(i, last, title, body)
        claims = _claims_from_blocks(i, text_blocks) + _claims_from_tables(i, tables)

        parsed_slides.append(ParsedSlide(
            index=i,
            slide_type=slide_type,
            text_blocks=text_blocks,
            tables=tables,
            images=images,
            claims=claims,
            title=title,
            body_items=body,
        ))

    return ParsedDeck(name=deck_name, slide_count=len(parsed_slides), slides=parsed_slides)
