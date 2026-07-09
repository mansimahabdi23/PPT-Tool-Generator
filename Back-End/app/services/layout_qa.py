"""Deterministic layout QA validator — no LLM.

Checks an output Presentation for geometric layout issues:
  · overflow  — shape extends past the right or bottom slide boundary
  · off_screen — shape starts at a negative position (left/top < 0)
  · empty_slide — slide has no visible text and no non-background images
  · zero_size — shape has width=0 or height=0 (usually a builder bug)

Geometry is math; all checks are pure arithmetic on EMU values.

Public API
----------
check(prs: Presentation) -> LayoutQAResult
"""

from __future__ import annotations

from dataclasses import dataclass, field

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from app import brand


# ---------------------------------------------------------------------------
# Result types
# ---------------------------------------------------------------------------


@dataclass
class LayoutIssue:
    """One layout defect."""

    slide_index: int
    shape_name: str
    issue_type: str  # "overflow" | "off_screen" | "empty_slide" | "zero_size"
    detail: str


@dataclass
class LayoutQAResult:
    passed: bool
    issues: list[LayoutIssue] = field(default_factory=list)
    slides_checked: int = 0


# ---------------------------------------------------------------------------
# Shape geometry checks
# ---------------------------------------------------------------------------


def _check_shape_geometry(
    shape: object,
    slide_index: int,
    slide_w: int,
    slide_h: int,
    issues: list[LayoutIssue],
) -> None:
    left: int = getattr(shape, "left", None) or 0
    top: int = getattr(shape, "top", None) or 0
    width: int = getattr(shape, "width", None) or 0
    height: int = getattr(shape, "height", None) or 0
    name: str = getattr(shape, "name", "") or ""

    if width == 0 or height == 0:
        issues.append(LayoutIssue(
            slide_index=slide_index,
            shape_name=name,
            issue_type="zero_size",
            detail=f"shape has zero dimension: width={width} height={height}",
        ))
        return  # no point checking overflow on a zero-size shape

    if left < 0 or top < 0:
        issues.append(LayoutIssue(
            slide_index=slide_index,
            shape_name=name,
            issue_type="off_screen",
            detail=f"shape starts off-screen: left={left} top={top}",
        ))

    right = left + width
    bottom = top + height
    # Allow 1% tolerance for rounding in EMU arithmetic
    overflow_x = right > slide_w * 1.01
    overflow_y = bottom > slide_h * 1.01

    if overflow_x or overflow_y:
        details = []
        if overflow_x:
            details.append(f"right={right} > slide_width={slide_w}")
        if overflow_y:
            details.append(f"bottom={bottom} > slide_height={slide_h}")
        issues.append(LayoutIssue(
            slide_index=slide_index,
            shape_name=name,
            issue_type="overflow",
            detail="; ".join(details),
        ))


# ---------------------------------------------------------------------------
# Empty-slide check
# ---------------------------------------------------------------------------


def _slide_has_content(slide: object, slide_w: int, slide_h: int) -> bool:
    """True if the slide has any text or any non-background image."""
    bg_w_thresh = slide_w * 0.90
    bg_h_thresh = slide_h * 0.90

    for shape in slide.shapes:  # type: ignore[attr-defined]
        # Non-empty text
        if getattr(shape, "has_text_frame", False):
            try:
                if shape.text_frame.text.strip():
                    return True
            except Exception:
                pass

        # Non-background picture (< 90% of slide dimensions)
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            w = getattr(shape, "width", 0) or 0
            h = getattr(shape, "height", 0) or 0
            if not (w >= bg_w_thresh and h >= bg_h_thresh):
                return True

    return False


# ---------------------------------------------------------------------------
# Shape tree traversal (one level of group recursion)
# ---------------------------------------------------------------------------


def _traverse(
    shapes: object,
    slide_index: int,
    slide_w: int,
    slide_h: int,
    issues: list[LayoutIssue],
    depth: int = 0,
) -> None:
    for shape in shapes:  # type: ignore[iteration]
        try:
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                if depth == 0:
                    _traverse(
                        shape.shapes, slide_index, slide_w, slide_h,
                        issues, depth=1,
                    )
                continue
            _check_shape_geometry(shape, slide_index, slide_w, slide_h, issues)
        except Exception:
            pass


# ---------------------------------------------------------------------------
# Public API
# ---------------------------------------------------------------------------


def check(prs: Presentation) -> LayoutQAResult:
    """Run layout geometry checks against the output *prs*.

    Checks every shape for overflow and off-screen placement, and every slide
    for empty content. Returns LayoutQAResult; passed=True iff issues is empty.
    """
    issues: list[LayoutIssue] = []
    slide_w: int = prs.slide_width
    slide_h: int = prs.slide_height

    for slide_index, slide in enumerate(prs.slides):
        _traverse(slide.shapes, slide_index, slide_w, slide_h, issues)

        if not _slide_has_content(slide, slide_w, slide_h):
            issues.append(LayoutIssue(
                slide_index=slide_index,
                shape_name="<slide>",
                issue_type="empty_slide",
                detail="slide has no visible text and no non-background images",
            ))

    return LayoutQAResult(
        passed=len(issues) == 0,
        issues=issues,
        slides_checked=len(prs.slides),
    )
