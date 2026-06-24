"""Infographic fragment merger (§5, §12).

Implements the XML-level shape-tree copy: deep-copies the shape children from a
source PPTX slide into the target slide, wrapped in a group shape whose transform
scales the entire source canvas into the target region.  Referenced image parts are
re-registered on the target slide so the result is self-contained.

No Aspose.Slides — this is purely python-pptx + lxml.

Public API
----------
merge_fragment(target_slide, src_pptx_path, src_slide_index, region)
"""

from __future__ import annotations

from copy import deepcopy
from pathlib import Path
from typing import Any

from lxml import etree  # type: ignore[import-untyped]
from pptx import Presentation
from pptx.oxml.ns import qn

from app.templates.layouts import Rect

# OOXML namespace URIs
_P = "http://schemas.openxmlformats.org/presentationml/2006/main"
_A = "http://schemas.openxmlformats.org/drawingml/2006/main"
_R = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"


# ---------------------------------------------------------------------------
# Internal helpers
# ---------------------------------------------------------------------------

def _next_shape_id(sp_tree: Any) -> int:
    """Return the next unused shape id in *sp_tree*."""
    ids = [int(el.get("id", 0)) for el in sp_tree.iter() if el.get("id") is not None]
    return (max(ids) if ids else 0) + 1


def _renumber_ids(grp_el: Any, start_id: int) -> int:
    """Reassign every ``id`` attribute in *grp_el* to avoid conflicts with the
    target slide.  Returns the next available id after renumbering."""
    counter = start_id
    for el in grp_el.iter():
        if el.get("id") is not None:
            el.set("id", str(counter))
            counter += 1
    return counter


def _copy_image_rels(src_slide: Any, target_slide: Any, grp_el: Any) -> None:
    """Rewrite r:embed / r:link attributes in *grp_el* so copied image shapes
    reference parts that exist in the target slide's package."""
    rId_map: dict[str, str] = {}

    embed_attr = f"{{{_R}}}embed"
    link_attr = f"{{{_R}}}link"

    # First pass: discover all rIds referenced in the copy and add their parts
    for el in grp_el.iter():
        for attr in (embed_attr, link_attr):
            old_rId: str | None = el.get(attr)
            if old_rId and old_rId not in rId_map:
                try:
                    rel = src_slide.part.rels[old_rId]
                    new_rId: str = target_slide.part.relate_to(
                        rel.target_part, rel.reltype
                    )
                    rId_map[old_rId] = new_rId
                except (KeyError, AttributeError):
                    pass  # unknown rel — leave as-is; PowerPoint will repair

    # Second pass: rewrite the rId values in the copied XML
    for el in grp_el.iter():
        for attr in (embed_attr, link_attr):
            old_rId = el.get(attr)
            if old_rId and old_rId in rId_map:
                el.set(attr, rId_map[old_rId])


def _build_group_wrapper(region: Rect, src_width: int, src_height: int, shape_id: int) -> Any:
    """Build a <p:grpSp> whose transform maps the source slide canvas to *region*.

    The child coordinate space (chOff / chExt) covers the full source slide so
    every source shape—regardless of its position—lands proportionally inside
    the target region.
    """
    grp_xml = (
        f'<p:grpSp xmlns:p="{_P}" xmlns:a="{_A}">'
        f"  <p:nvGrpSpPr>"
        f'    <p:cNvPr id="{shape_id}" name="InfographicFragment"/>'
        f"    <p:cNvGrpSpPr/>"
        f"    <p:nvPr/>"
        f"  </p:nvGrpSpPr>"
        f"  <p:grpSpPr>"
        f"    <a:xfrm>"
        f'      <a:off x="{region.left}" y="{region.top}"/>'
        f'      <a:ext cx="{region.width}" cy="{region.height}"/>'
        f'      <a:chOff x="0" y="0"/>'
        f'      <a:chExt cx="{src_width}" cy="{src_height}"/>'
        f"    </a:xfrm>"
        f"  </p:grpSpPr>"
        f"</p:grpSp>"
    )
    return etree.fromstring(grp_xml)


# ---------------------------------------------------------------------------
# Public API
# ---------------------------------------------------------------------------

def merge_fragment(
    target_slide: Any,
    src_pptx_path: Path,
    src_slide_index: int,
    region: Rect,
) -> None:
    """Deep-copy shapes from *src_pptx_path*[*src_slide_index*] into *target_slide*.

    Shapes are wrapped in a group whose transform scales the full source slide
    canvas (e.g. 12 192 000 × 6 858 000 EMU) to *region* on the target slide.
    Referenced image parts are copied so the output is self-contained.

    Parameters
    ----------
    target_slide    : python-pptx Slide object (the slide to merge into)
    src_pptx_path   : path to the source .pptx file
    src_slide_index : 0-based index of the slide to pull shapes from
    region          : target Rect on the target_slide in EMU
    """
    src_prs: Any = Presentation(str(src_pptx_path))
    src_slide: Any = src_prs.slides[src_slide_index]

    src_width: int = int(src_prs.slide_width)
    src_height: int = int(src_prs.slide_height)

    src_cSld: Any = src_slide.element.find(qn("p:cSld"))
    src_spTree: Any = src_cSld.find(qn("p:spTree"))

    target_spTree: Any = target_slide.shapes._spTree

    # Reserve a shape ID range for the new group (will be renumbered below)
    start_id = _next_shape_id(target_spTree)

    # Build the wrapper group shape
    grp_el = _build_group_wrapper(region, src_width, src_height, start_id)

    # Copy non-structural children from source spTree into the group
    _skip_tags = {qn("p:nvGrpSpPr"), qn("p:grpSpPr")}
    for child in src_spTree:
        if child.tag not in _skip_tags:
            grp_el.append(deepcopy(child))

    # Renumber IDs to avoid conflicts with the target slide
    _renumber_ids(grp_el, start_id)

    # Copy image parts (updates r:embed / r:link in grp_el in-place)
    _copy_image_rels(src_slide, target_slide, grp_el)

    # Append the group to the target slide
    target_spTree.append(grp_el)
