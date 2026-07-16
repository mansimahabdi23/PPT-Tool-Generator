"""Clone-and-fill composer for body-block, cover, and closing slides.

Strategy (body-block — template slide 5)
-----------------------------------------
1. Deep-copy all shapes from template slide 5 into a new blank slide,
   SKIPPING "Rounded Rectangle 1" (the layout guide box — not a rendered element).
2. Register image part relationships (logo, etc.) in the output package.
3. Replace the title run text in TextBox 5, preserving run formatting.
4. Inject a new body textbox inside the guide box's bounding box region.
5. Overflow: use normAutofit so PowerPoint/LibreOffice shrinks if needed;
   only flag slides that are genuinely extreme (> _EXTREME_OVERFLOW items).

Strategy (cover — template slide 1)
-------------------------------------
Clone template slide 1 intact (photo, overlay, logo preserved) then
replace text in-place:
  TextBox 7  → source title (one run; extra paragraphs removed)
  TextBox 9  → source body_items as paragraphs (preserves run formatting)
  TextBox 12 ("Deck Name") is left unchanged.

Strategy (closing — template slide 20)
-----------------------------------------
Clone template slide 20 exactly as-is (fixed brand end-card, no text
replacement).  Append as a new final slide after all content slides.

Public API
----------
load_template(path)                                       -> Presentation
clone_body_block_slide(template_prs, out_prs, title,
                       body_items, slide_number)          -> bool
clone_cover_slide(template_prs, out_prs, title,
                  body_items)                             -> None
clone_closing_slide(template_prs, out_prs)                -> None
"""

from __future__ import annotations

import logging
from copy import deepcopy
from pathlib import Path
from typing import Any

from lxml import etree  # type: ignore[import-untyped]
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from pptx.util import Emu, Pt

from app import brand
from app.theme import LIGHT, ThemePalette

logger = logging.getLogger(__name__)

# ---------------------------------------------------------------------------
# Constants
# ---------------------------------------------------------------------------

# Template slide indices (0-based)
_TEMPLATE_SLIDE_IDX   = 4   # slide 5 — body-block
_TEMPLATE_COVER_IDX   = 0   # slide 1 — dark cover with photo/graphic/logo
_TEMPLATE_CLOSING_IDX = 19  # slide 20 — "Skills Visibility. Business Agility."

# Shape names in template slide 1 (cover)
_COVER_HEADING_SHAPE = "TextBox 7"   # two paragraphs → "<Cover Slide" / "Heading Text>"
_COVER_BODY_SHAPE    = "TextBox 9"   # one paragraph  → "<Cover Slide subheading…>"
# TextBox 12 ("Deck Name") is intentionally left unchanged

# Cover bullet zone — placed below TextBox 9 (bot=4,397,372 EMU), left of photo (5.85 in)
# Exact EMU from template: TextBox 9 left=595,313 top=3,689,486 h=707,886 → bot=4,397,372
_COVER_BULLETS_LEFT   = 595_311    # 0.65 in — flush with TextBox 9
_COVER_BULLETS_TOP    = 4_488_812  # 0.10 in below TextBox 9 bottom
_COVER_BULLETS_WIDTH  = 4_326_311  # 4.73 in — same width as TextBox 9, left of photo
_COVER_BULLETS_HEIGHT = 1_188_720  # 1.30 in — to ~6.20 in (TextBox 12 starts at 6.38 in)
_COVER_BULLETS_FONT_PT = 11

# OOXML namespace URIs
_A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
_R    = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"

# Body box geometry: Rounded Rectangle 1 bounds with 0.2in (182 880 EMU) inset
#   Rounded Rectangle 1: x=407 306, y=976 393, cx=11 375 996, cy=5 131 799
_PAD = 182_880
_BODY_LEFT   = 407_306  + _PAD          #  590 186 EMU
_BODY_TOP    = 976_393  + _PAD          # 1 159 273 EMU
_BODY_WIDTH  = 11_375_996 - 2 * _PAD   # 11 010 236 EMU
_BODY_HEIGHT = 5_131_799  - 2 * _PAD   #  4 766 039 EMU

# Body font: fixed 12 pt Poppins; normAutofit handles any real overflow in PowerPoint.
_BODY_FONT_PT = 12

# Spacing: 140 % line spacing; 8 pt space_before between bullets.
_LINE_SPACING_PCT = 140_000   # OOXML spcPct val: 140 000 = 140%
_PARA_SPACE_PT    = 8

# Flag slides only when content is genuinely extreme (no layout can hold it).
# Matches _OVERFLOW_FLAG_THRESHOLD in layout_catalog.py.
_EXTREME_OVERFLOW = 12

# Text-body insets (EMU): 0.1 in on all sides; combined with the 0.2 in box
# padding this gives the text ~0.3 in clearance from the rounded-rect edge.
_INSET_EMU = 91_440   # 0.1 in = 91 440 EMU

# Icon accent — upper-right corner of Rounded Rectangle 1.
# 1.1 in square with a white backing tile so the lavender icon reads clearly
# against the lavender box background.
#   Rounded Rectangle 1: x=407 306, y=976 393, cx=11 375 996, cy=5 131 799
_RRECT_TOP    = 976_393
_RRECT_RIGHT  = 407_306 + 11_375_996   # 11 783 302 EMU (right edge of rect)
_ICON_SIZE    = 1_005_840               # 1.10 in square
_ICON_MARGIN  = 228_600                 # 0.25 in inset from the rect edges
_ICON_LEFT    = _RRECT_RIGHT - _ICON_MARGIN - _ICON_SIZE   # 10 548 862 EMU
_ICON_TOP     = _RRECT_TOP   + _ICON_MARGIN                #  1 204 993 EMU

# White backing tile — adds 0.08 in breathing room around the icon
_ICON_BG_PAD  = 73_152                                      # 0.08 in
_ICON_BG_SIZE = _ICON_SIZE + 2 * _ICON_BG_PAD              # ~1.26 in square
_ICON_BG_LEFT = _ICON_LEFT - _ICON_BG_PAD
_ICON_BG_TOP  = _ICON_TOP  - _ICON_BG_PAD

# spTree structural tags — never copied or removed by mistake
_STRUCT_TAGS = {qn("p:nvGrpSpPr"), qn("p:grpSpPr")}

# Logo picture on body-block slides sits below 6.5" — used to find it by y-position
# so we can swap the dark logo → logo_white.png when rendering dark-theme slides.
_LOGO_Y_MIN_EMU = int(6.5 * 914_400)   # ~5 943 600 EMU

# Name of the light-purple rounded-rectangle guide box on body-block slides.
# This shape defines the layout region but must NOT be rendered in the output.
_GUIDE_BOX_NAME = "Rounded Rectangle 1"

# ---------------------------------------------------------------------------
# Infographic slide — Bookmark Card fragment (5-slot parallel capabilities)
# ---------------------------------------------------------------------------

# Path relative to assets_root
_BOOKMARK_CARD_FRAG = (
    "infographics/Bookmark Card Process Flow – 5 Stage Fading Cards (Purple).pptx"
)

# ---------------------------------------------------------------------------
# Infographic slide — Hub-and-Spoke fragment (4-slot parallel features)
# ---------------------------------------------------------------------------

_HUB_SPOKE_FRAG = (
    "infographics/4-Feature Hub-and-Spoke Diagram – Central Brand Icon (Purple).pptx"
)
_HUB_SPOKE_MAX_ITEMS = 4
# "Feature" label boxes beside each satellite node — too narrow (76pt) for our items; clear them
_HUB_HEADLINE_LABEL = "Text Placeholder 33"
# Wide body-bar boxes (254pt) beside each node — filled with item text
_HUB_BODY_LABEL = "Rectangle 69"

# Oval badge shapes to remove — numbered 1–5 circles misrepresent parallel items
_CARD_BADGE_NAMES: frozenset[str] = frozenset({
    "Oval 11", "Oval 14", "Oval 17", "Oval 20", "Oval 23",
})

# Headline slots in left-to-right card order
_CARD_HEADLINE_NAMES: tuple[str, ...] = (
    "TextBox 12", "TextBox 15", "TextBox 18", "TextBox 21", "TextBox 24",
)

# Body slots — lorem ipsum placeholder text to be cleared
_CARD_BODY_NAMES: frozenset[str] = frozenset({
    "TextBox 13", "TextBox 16", "TextBox 19", "TextBox 22", "TextBox 25",
})

# Fragment's own category/title header — removed after merge (body-block has its own title)
_FRAG_HEADER_NAMES: frozenset[str] = frozenset({"TextBox 2", "TextBox 3"})

# Card background shapes — their fill is overridden to uniform brand-purple so all 5
# are equally readable (the fragment uses a fading-opacity design that makes cards 4-5
# near-invisible with white text).
_CARD_FILL_NAMES: frozenset[str] = frozenset({
    "Freeform 6",
    "Rounded Rectangle 7",
    "Rounded Rectangle 8",
    "Rounded Rectangle 9",
    "Rounded Rectangle 10",
})

# Bookmark Card fragment: maximum items the 5-slot fragment can hold
_BOOKMARK_CARD_MAX_ITEMS = 5

# Card 1 height fix: Freeform 6 (the bookmark shape) has a pointed tail that makes it
# much taller (~271 pt) than the 4 rounded-rectangle cards (~158 pt).  Clip to match.
_CARD_ROUNDED_H_EMU = 2_006_600  # 158 pt × 12 700 EMU/pt — height of cards 2-5


# ---------------------------------------------------------------------------
# Private helpers
# ---------------------------------------------------------------------------

def _next_shape_id(sp_tree: Any) -> int:
    """Return the next unused shape id (max existing + 1)."""
    ids = [int(el.get("id", 0)) for el in sp_tree.iter() if el.get("id") is not None]
    return (max(ids) if ids else 0) + 1


def _register_image_rels(src_slide: Any, target_slide: Any, tree_el: Any) -> None:
    """Rewrite r:embed / r:link attributes in *tree_el* to reference parts
    that exist in *target_slide*'s package.  Same two-pass strategy as merge.py."""
    embed_attr = f"{{{_R}}}embed"
    link_attr  = f"{{{_R}}}link"
    rId_map: dict[str, str] = {}

    for el in tree_el.iter():
        for attr in (embed_attr, link_attr):
            old = el.get(attr)
            if old and old not in rId_map:
                try:
                    rel = src_slide.part.rels[old]
                    rId_map[old] = target_slide.part.relate_to(rel.target_part, rel.reltype)
                except (KeyError, AttributeError):
                    pass  # unknown rel — PowerPoint will repair

    for el in tree_el.iter():
        for attr in (embed_attr, link_attr):
            old = el.get(attr)
            if old and old in rId_map:
                el.set(attr, rId_map[old])


def _replace_title(sp_tree: Any, new_title: str) -> bool:
    """Find 'TextBox 5' in *sp_tree* and replace its first run text.

    Edits the <a:r><a:t> node directly (not text_frame.text) so all run-level
    formatting (font, size, color) is preserved from the template.  Any extra
    runs in the same paragraph are removed to avoid leftover placeholder text.

    Returns True on success.
    """
    for sp in sp_tree.findall(f'.//{qn("p:sp")}'):
        cNvPr = sp.find(f'.//{qn("p:cNvPr")}')
        if cNvPr is None or cNvPr.get("name") != "TextBox 5":
            continue
        runs = sp.findall(f'.//{qn("a:r")}')
        if not runs:
            continue
        t_el = runs[0].find(qn("a:t"))
        if t_el is not None:
            t_el.text = new_title
        # Remove any extra runs in the same paragraph
        first_para = runs[0].getparent()
        for extra in runs[1:]:
            if extra.getparent() is first_para:
                first_para.remove(extra)
        return True
    return False


def _update_slide_number(sp_tree: Any, slide_number: int) -> None:
    """Update the cached display text of the <a:fld type="slidenum"> element."""
    for fld in sp_tree.findall(f'.//{qn("a:fld")}'):
        if fld.get("type") == "slidenum":
            t_el = fld.find(qn("a:t"))
            if t_el is not None:
                t_el.text = str(slide_number)
            return


def _style_body_frame(tf: Any) -> None:
    """Apply vertical centering, insets, and normAutofit to a text frame.

    normAutofit tells PowerPoint to shrink the font automatically if the text
    exceeds the box — a reliable native behaviour that replaces the previous
    Python estimation cascade.  The box size stays fixed so the layout is
    always preserved.
    """
    txBody = tf._txBody
    bodyPr = txBody.find(qn("a:bodyPr"))
    if bodyPr is None:
        bodyPr = etree.SubElement(txBody, qn("a:bodyPr"))

    # Vertical centering — text floats in the middle of the 5.21" tall box
    bodyPr.set("anchor", "ctr")
    # Fixed insets (0.1 in all sides; combined with 0.2 in box padding = 0.3 in clearance)
    bodyPr.set("lIns", str(_INSET_EMU))
    bodyPr.set("rIns", str(_INSET_EMU))
    bodyPr.set("tIns", str(_INSET_EMU))
    bodyPr.set("bIns", str(_INSET_EMU))

    # Replace any existing autofit element with normAutofit so PowerPoint
    # shrinks the font natively when the text genuinely overflows the box.
    for tag in (qn("a:noAutofit"), qn("a:normAutofit"), qn("a:spAutoFit")):
        for el in bodyPr.findall(tag):
            bodyPr.remove(el)
    bodyPr.append(
        etree.fromstring(
            f'<a:normAutofit xmlns:a="{_A_NS}"/>'
        )
    )


def _set_para_spacing(para: Any, is_first: bool) -> None:
    """Set 140 % line spacing and paragraph space_before on *para*."""
    para.space_before = Pt(0) if is_first else Pt(_PARA_SPACE_PT)

    pPr = para._p.get_or_add_pPr()
    for old in pPr.findall(qn("a:lnSpc")):
        pPr.remove(old)
    pPr.append(
        etree.fromstring(
            f'<a:lnSpc xmlns:a="{_A_NS}">'
            f'<a:spcPct val="{_LINE_SPACING_PCT}"/>'
            f"</a:lnSpc>"
        )
    )


def _inject_body(slide: Any, body_items: list[str], theme: ThemePalette) -> bool:
    """Add a body textbox to *slide* at the body-box geometry.

    Uses a fixed 12 pt Poppins font with normAutofit so PowerPoint/LibreOffice
    shrinks the text automatically if it genuinely overflows.  Only flags a
    slide when the item count is extreme (> _EXTREME_OVERFLOW), matching the
    plan-stage threshold so the signal is coherent end-to-end.

    Returns True when the slide is flagged for human review.
    """
    overflow_flagged = len(body_items) > _EXTREME_OVERFLOW
    if overflow_flagged:
        logger.warning(
            "clone_body_block: %d items exceeds threshold %d — slide flagged",
            len(body_items), _EXTREME_OVERFLOW,
        )

    txBox = slide.shapes.add_textbox(
        Emu(_BODY_LEFT), Emu(_BODY_TOP), Emu(_BODY_WIDTH), Emu(_BODY_HEIGHT)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    _style_body_frame(tf)

    r, g, b = brand.rgb_tuple(theme.body_text)
    for i, item in enumerate(body_items):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        _set_para_spacing(para, is_first=(i == 0))
        run = para.add_run()
        run.text = item
        run.font.name  = brand.POPPINS
        run.font.size  = Pt(_BODY_FONT_PT)
        run.font.color.rgb = RGBColor(r, g, b)

    return overflow_flagged


def _inject_icon(slide: Any, icon_path: Path, theme: ThemePalette) -> None:
    """Place a 1.1 in icon accent in the upper-right of the body block.

    Renders a rounded-rectangle backing (color from *theme*) then the PNG on
    top.  Failures are logged and swallowed — a missing file never breaks the
    slide.
    """
    try:
        # Backing tile — color adapts to theme (white on light, dark-purple on dark)
        bg = slide.shapes.add_shape(
            5,  # MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE
            Emu(_ICON_BG_LEFT),
            Emu(_ICON_BG_TOP),
            Emu(_ICON_BG_SIZE),
            Emu(_ICON_BG_SIZE),
        )
        bg.fill.solid()
        ir, ig, ib = brand.rgb_tuple(theme.icon_tile)
        bg.fill.fore_color.rgb = RGBColor(ir, ig, ib)
        # Remove default PowerPoint border — white tile should be borderless
        spPr = bg.element.find(qn("p:spPr"))
        if spPr is not None:
            for ln in spPr.findall(qn("a:ln")):
                spPr.remove(ln)
            spPr.append(etree.fromstring(
                f'<a:ln xmlns:a="{_A_NS}"><a:noFill/></a:ln>'
            ))

        # Icon PNG on top of the white backing
        slide.shapes.add_picture(
            str(icon_path),
            Emu(_ICON_LEFT),
            Emu(_ICON_TOP),
            Emu(_ICON_SIZE),
            Emu(_ICON_SIZE),
        )
    except Exception:
        logger.warning(
            "clone_body_block: failed to insert icon %s — slide renders without it",
            icon_path,
            exc_info=True,
        )


# ---------------------------------------------------------------------------
# Cover/closing helpers
# ---------------------------------------------------------------------------

def _shape_name(child: Any) -> str | None:
    """Return the cNvPr name attribute of *child*, or None if not present."""
    cNvPr = child.find(f'.//{qn("p:cNvPr")}')
    return cNvPr.get("name") if cNvPr is not None else None


def _replace_cover_heading(sp_tree: Any, title: str) -> bool:
    """Write *title* into TextBox 7 (the cover heading shape).

    Template slide 1 stores the placeholder text across **two paragraphs**:
    paragraph 0 = '<Cover Slide' and paragraph 1 = 'Heading Text>'.
    We put the entire title into paragraph 0's first run and remove paragraph 1,
    preserving the run's font/size/colour from the template.

    Returns True on success.
    """
    for sp in sp_tree.findall(f'.//{qn("p:sp")}'):
        cNvPr = sp.find(f'.//{qn("p:cNvPr")}')
        if cNvPr is None or cNvPr.get("name") != _COVER_HEADING_SHAPE:
            continue
        txBody = sp.find(qn("p:txBody"))
        if txBody is None:
            return False
        paras = txBody.findall(qn("a:p"))
        if not paras:
            return False
        # Replace first run in paragraph 0 with the title text
        runs = paras[0].findall(qn("a:r"))
        if runs:
            t = runs[0].find(qn("a:t"))
            if t is not None:
                t.text = title
            for extra in runs[1:]:
                paras[0].remove(extra)
        # Remove extra paragraphs (paragraph 1 = 'Heading Text>')
        for extra_p in paras[1:]:
            txBody.remove(extra_p)
        # Add normAutofit so PowerPoint shrinks heading to fit within the 1.58 in box
        bodyPr = txBody.find(qn("a:bodyPr"))
        if bodyPr is not None:
            for tag in (qn("a:noAutofit"), qn("a:normAutofit"), qn("a:spAutoFit")):
                for el in bodyPr.findall(tag):
                    bodyPr.remove(el)
            bodyPr.append(etree.fromstring(f'<a:normAutofit xmlns:a="{_A_NS}"/>'))
        return True
    return False


def _replace_cover_body(sp_tree: Any, body_items: list[str]) -> bool:
    """Write *body_items* into TextBox 9 (cover subheading/description shape).

    Paragraph 0 receives body_items[0] (the subtitle/subheading).
    Each subsequent item becomes a new paragraph that inherits the first
    paragraph's run formatting (font, size, colour) from the template.

    Returns True on success.
    """
    for sp in sp_tree.findall(f'.//{qn("p:sp")}'):
        cNvPr = sp.find(f'.//{qn("p:cNvPr")}')
        if cNvPr is None or cNvPr.get("name") != _COVER_BODY_SHAPE:
            continue
        txBody = sp.find(qn("p:txBody"))
        if txBody is None:
            return False
        paras = txBody.findall(qn("a:p"))
        if not paras:
            return False
        # Clear any extra existing paragraphs (template has only one)
        for extra_p in paras[1:]:
            txBody.remove(extra_p)

        if not body_items:
            # Nothing to write — leave the placeholder cleared
            runs = paras[0].findall(qn("a:r"))
            for r in runs:
                t = r.find(qn("a:t"))
                if t is not None:
                    t.text = ""
            return True

        # Write body_items[0] (subtitle only) into paragraph 0.
        # Actual bullets go into a separate injected textbox — see _inject_cover_bullets().
        first_runs = paras[0].findall(qn("a:r"))
        if first_runs:
            t = first_runs[0].find(qn("a:t"))
            if t is not None:
                t.text = body_items[0]
            for extra in first_runs[1:]:
                paras[0].remove(extra)

        return True
    return False


def _inject_cover_bullets(slide: Any, bullet_items: list[str]) -> None:
    """Inject a small-text bullet zone below TextBox 9 on the cover slide.

    Constrained to the left half of the slide (≤ 5.20 in) to avoid overlapping
    the photo on the right.  11 pt Poppins white, normAutofit so the box shrinks
    if bullets are numerous.  Called only when body_items has more than one entry.
    """
    if not bullet_items:
        return

    txBox = slide.shapes.add_textbox(
        Emu(_COVER_BULLETS_LEFT),
        Emu(_COVER_BULLETS_TOP),
        Emu(_COVER_BULLETS_WIDTH),
        Emu(_COVER_BULLETS_HEIGHT),
    )
    tf = txBox.text_frame
    tf.word_wrap = True

    txBody = tf._txBody
    bodyPr = txBody.find(qn("a:bodyPr"))
    if bodyPr is None:
        bodyPr = etree.SubElement(txBody, qn("a:bodyPr"))
    for tag in (qn("a:noAutofit"), qn("a:normAutofit"), qn("a:spAutoFit")):
        for el in bodyPr.findall(tag):
            bodyPr.remove(el)
    bodyPr.append(etree.fromstring(f'<a:normAutofit xmlns:a="{_A_NS}"/>'))

    for i, item in enumerate(bullet_items):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if i > 0:
            para.space_before = Pt(4)
        run = para.add_run()
        run.text = item
        run.font.name = brand.POPPINS
        run.font.size = Pt(_COVER_BULLETS_FONT_PT)
        run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)


def _clone_template_slide(
    template_prs: Any,
    out_prs: Any,
    template_idx: int,
) -> Any:
    """Deep-copy all shapes from template slide *template_idx* into a new blank slide
    in *out_prs*.  Image relationships are re-registered.  Returns the new slide."""
    tmpl_slide: Any = template_prs.slides[template_idx]
    blank_layout = out_prs.slide_layouts[6]
    new_slide: Any = out_prs.slides.add_slide(blank_layout)

    src_spTree = tmpl_slide.element.find(qn("p:cSld")).find(qn("p:spTree"))
    tgt_spTree = new_slide.shapes._spTree

    for child in list(tgt_spTree):
        if child.tag not in _STRUCT_TAGS:
            tgt_spTree.remove(child)

    start_id = _next_shape_id(tgt_spTree)
    counter  = start_id
    for child in src_spTree:
        if child.tag in _STRUCT_TAGS:
            continue
        copied = deepcopy(child)
        for el in copied.iter():
            if el.get("id") is not None:
                el.set("id", str(counter))
                counter += 1
        tgt_spTree.append(copied)

    _register_image_rels(tmpl_slide, new_slide, tgt_spTree)
    return new_slide


# ---------------------------------------------------------------------------
# Infographic slide helpers
# ---------------------------------------------------------------------------

def _set_first_run_text(sp_el: Any, text: str) -> None:
    """Replace all run text in *sp_el* with *text*, keeping first run's formatting."""
    all_runs = sp_el.findall(f".//{qn('a:r')}")
    if not all_runs:
        return
    t_el = all_runs[0].find(qn("a:t"))
    if t_el is not None:
        t_el.text = text
    first_para = all_runs[0].getparent()
    for run in all_runs[1:]:
        if run.getparent() is first_para:
            first_para.remove(run)
    txBody = sp_el.find(f".//{qn('a:txBody')}")
    if txBody is not None:
        for p in txBody.findall(qn("a:p"))[1:]:
            txBody.remove(p)


def _clear_all_runs(sp_el: Any) -> None:
    """Remove every run element from *sp_el*, leaving the text box empty."""
    for r in list(sp_el.findall(f".//{qn('a:r')}")):
        parent = r.getparent()
        if parent is not None:
            parent.remove(r)


def _set_shape_solid_fill(sp_el: Any, color_hex: str) -> None:
    """Replace existing fill on *sp_el* with a uniform solidFill of *color_hex*.

    Strips any gradFill, noFill, pattFill, or blipFill already present so the
    result is a single unconditional solid color.  *color_hex* is 6-char hex
    without the leading '#' (matches brand.py convention).
    """
    spPr = sp_el.find(qn("p:spPr"))
    if spPr is None:
        return
    fill_tags = {
        qn("a:noFill"), qn("a:solidFill"), qn("a:gradFill"),
        qn("a:pattFill"), qn("a:blipFill"), qn("a:grpFill"),
    }
    for child in [c for c in list(spPr) if c.tag in fill_tags]:
        spPr.remove(child)
    solid_el = etree.fromstring(
        f'<a:solidFill xmlns:a="{_A_NS}">'
        f'<a:srgbClr val="{color_hex}"/>'
        f'</a:solidFill>'
    )
    # Insert fill after the last geometry element (xfrm / prstGeom / custGeom)
    insert_at = 0
    for i, child in enumerate(spPr):
        if child.tag in (qn("a:xfrm"), qn("a:prstGeom"), qn("a:custGeom")):
            insert_at = i + 1
    spPr.insert(insert_at, solid_el)


def _fix_headline_wrapping(sp_el: Any) -> None:
    """Prevent mid-Latin-word line breaks and widen the headline textbox.

    Sets latinLnBrk="0" on every paragraph so words are never split mid-character
    (fixes LibreOffice hyphenating e.g. "Inferenc-e Engine").  Also widens the
    textbox by 1.5× and re-centres it so long labels wrap at word boundaries
    with a comfortable column width.
    """
    # ---- 1. Disable Latin mid-word line breaks on every paragraph ----
    for para in sp_el.findall(f".//{qn('a:p')}"):
        pPr = para.find(qn("a:pPr"))
        if pPr is None:
            pPr = etree.Element(qn("a:pPr"))
            para.insert(0, pPr)
        pPr.set("latinLnBrk", "0")

    # ---- 2. Reduce font size so no single word exceeds the column width ----
    # 10 pt in a 1.5× widened box (~111 pt) allows ~18 chars per line — enough
    # for "Recommendations" (15) and "Conversational" (14) without any mid-word break.
    for rPr in sp_el.findall(f".//{qn('a:rPr')}"):
        rPr.set("sz", "1000")   # 1000 hundredths-of-a-point = 10 pt

    # ---- 3. Widen the textbox geometry (1.5×) and re-centre it ----
    spPr = sp_el.find(qn("p:spPr"))
    if spPr is None:
        return
    xfrm = spPr.find(qn("a:xfrm"))
    if xfrm is None:
        return
    off = xfrm.find(qn("a:off"))
    ext = xfrm.find(qn("a:ext"))
    if off is None or ext is None:
        return
    try:
        old_cx = int(ext.get("cx", 0))
        old_x  = int(off.get("x", 0))
        new_cx = int(old_cx * 1.5)
        # Shift left by half the extra width to keep the box centred over its card
        new_x  = old_x - (new_cx - old_cx) // 2
        ext.set("cx", str(new_cx))
        off.set("x",  str(new_x))
    except (ValueError, TypeError):
        pass  # malformed geometry — leave untouched


def _fix_card1_bookmark_tail(sp_el: Any) -> None:
    """Convert Freeform 6 from a bookmark (pointed tail) to a roundRect, clipping height.

    The bookmark freeform is ~271 pt tall — more than the ~158 pt rounded-rectangle
    cards 2-5.  Replacing the custom geometry with a standard roundRect and capping
    the height at _CARD_ROUNDED_H_EMU makes all 5 cards the same visual height.
    """
    spPr = sp_el.find(qn("p:spPr"))
    if spPr is None:
        return

    # Clip height to match the rounded-rectangle cards
    xfrm = spPr.find(qn("a:xfrm"))
    if xfrm is not None:
        ext = xfrm.find(qn("a:ext"))
        if ext is not None:
            ext.set("cy", str(_CARD_ROUNDED_H_EMU))

    # Replace custom bookmark geometry with a standard roundRect
    cust = spPr.find(qn("a:custGeom"))
    if cust is not None:
        spPr.remove(cust)
        insert_at = 0
        for i, child in enumerate(spPr):
            if child.tag == qn("a:xfrm"):
                insert_at = i + 1
                break
        spPr.insert(insert_at, etree.fromstring(
            f'<a:prstGeom xmlns:a="{_A_NS}" prst="roundRect">'
            f'<a:avLst/>'
            f'</a:prstGeom>'
        ))


def _fill_bookmark_card_group(grp_el: Any, items: list[str]) -> None:
    """Post-merge cleanup for the Bookmark Card merged group.

    1. Remove numbered oval badges (items are parallel, not sequential).
    2. Remove fragment's own header textboxes (body-block chrome has the title).
    3. Fill each Headline slot with the corresponding item text.
    4. Clear each Body slot (removes lorem ipsum).
    """
    headline_map: dict[str, str] = {
        name: items[i]
        for i, name in enumerate(_CARD_HEADLINE_NAMES)
        if i < len(items)
    }
    remove_names = _CARD_BADGE_NAMES | _FRAG_HEADER_NAMES

    for child in list(grp_el):
        name = _shape_name(child)
        if name in remove_names:
            grp_el.remove(child)
        elif name in _CARD_FILL_NAMES:
            _set_shape_solid_fill(child, brand.BRAND_PURPLE)
            if name == "Freeform 6":
                _fix_card1_bookmark_tail(child)
        elif name in headline_map:
            _set_first_run_text(child, headline_map[name])
            _fix_headline_wrapping(child)
        elif name in _CARD_BODY_NAMES:
            _clear_all_runs(child)


def _get_sp_xy(sp_el: Any) -> tuple[int, int]:
    """Return (x, y) offset of *sp_el* from its xfrm element (for position-based sorting)."""
    xfrm = sp_el.find(f".//{qn('a:xfrm')}")
    if xfrm is None:
        return (0, 0)
    off = xfrm.find(qn("a:off"))
    if off is None:
        return (0, 0)
    return (int(off.get("x", 0)), int(off.get("y", 0)))


def _fill_hub_spoke_group(grp_el: Any, items: list[str]) -> None:
    """Post-merge cleanup for the Hub-and-Spoke merged group.

    1. Remove fragment header textboxes (TextBox 2 / 3).
    2. Fill each wide body bar (Rectangle 69 × 4) with item text, sorted
       left-to-right (x) then top-to-bottom (y) within each side.
    3. Clear the narrow "Feature" headline labels (Text Placeholder 33 × 4)
       — they are too small (76 pt wide) to hold our full-sentence items.
    """
    # Remove fragment header shapes from the merged wrapper's direct children
    for child in list(grp_el):
        if _shape_name(child) in _FRAG_HEADER_NAMES:
            grp_el.remove(child)

    # Collect body bars and headline labels from anywhere inside the merged group
    body_els: list[Any] = []
    headline_els: list[Any] = []
    for sp in grp_el.iter(qn("p:sp")):
        n = _shape_name(sp)
        if n == _HUB_BODY_LABEL:
            body_els.append(sp)
        elif n == _HUB_HEADLINE_LABEL:
            headline_els.append(sp)

    # Sort body bars by (x, y): left column first (lower x), top before bottom within each column
    body_els.sort(key=_get_sp_xy)

    # Fill body bars — preserve existing run formatting (white text on dark bg)
    for i, sp_el in enumerate(body_els):
        if i >= len(items):
            break
        _set_first_run_text(sp_el, items[i])
        # Prevent mid-word breaks
        for para in sp_el.findall(f".//{qn('a:p')}"):
            pPr = para.find(qn("a:pPr"))
            if pPr is None:
                pPr = etree.Element(qn("a:pPr"))
                para.insert(0, pPr)
            pPr.set("latinLnBrk", "0")
        # Reduce to 9 pt so sentences wrap within the ~24 pt tall bar
        for rPr in sp_el.findall(f".//{qn('a:rPr')}"):
            rPr.set("sz", "900")
        # normAutofit: PowerPoint shrinks further if any line still overflows
        txBody = sp_el.find(f".//{qn('a:txBody')}")
        if txBody is not None:
            bodyPr = txBody.find(qn("a:bodyPr"))
            if bodyPr is None:
                bodyPr = etree.SubElement(txBody, qn("a:bodyPr"))
            for tag in (qn("a:noAutofit"), qn("a:normAutofit"), qn("a:spAutoFit")):
                for el in bodyPr.findall(tag):
                    bodyPr.remove(el)
            bodyPr.append(etree.fromstring(f'<a:normAutofit xmlns:a="{_A_NS}"/>'))

    # Clear "Feature" headline labels — not used for our full-sentence items
    for sp_el in headline_els:
        _clear_all_runs(sp_el)


# ---------------------------------------------------------------------------
# Dark-theme logo helper
# ---------------------------------------------------------------------------

def _swap_logo_dark(new_slide: Any, tgt_spTree: Any, assets_root: Path) -> None:
    """Replace the dark-on-light logo with logo_white.png for dark-theme slides.

    Finds the logo ``<p:pic>`` element (identified by y-position > 6.5"),
    reads its position/size from the transform, removes it, then adds
    ``logo_white.png`` from *assets_root*/logo/ at the exact same geometry
    using python-pptx's standard add_picture().
    """
    logo_white = assets_root / "logo" / "logo_white.png"
    if not logo_white.exists():
        logger.warning(
            "_swap_logo_dark: logo_white.png not found at %s — logo not swapped",
            logo_white,
        )
        return

    for pic_el in list(tgt_spTree.findall(qn("p:pic"))):
        spPr = pic_el.find(qn("p:spPr"))
        xfrm = spPr.find(qn("a:xfrm")) if spPr is not None else None
        off  = xfrm.find(qn("a:off")) if xfrm is not None else None
        ext  = xfrm.find(qn("a:ext")) if xfrm is not None else None
        if off is None or ext is None:
            continue
        y = int(off.get("y", 0))
        if y < _LOGO_Y_MIN_EMU:
            continue  # not the logo

        x  = int(off.get("x", 0))
        cx = int(ext.get("cx", 0))
        cy = int(ext.get("cy", 0))
        parent = pic_el.getparent()
        if parent is not None:
            parent.remove(pic_el)

        new_slide.shapes.add_picture(str(logo_white), Emu(x), Emu(y), Emu(cx), Emu(cy))
        return


# ---------------------------------------------------------------------------
# Public API
# ---------------------------------------------------------------------------

def load_template(path: Path) -> Any:
    """Open the iMocha template PPTX and return a python-pptx Presentation.

    Call once per compose run and pass the result to clone_body_block_slide
    for every body-block slide to avoid repeated file I/O.
    """
    return Presentation(str(path))


def clone_infographic_slide(
    template_prs: Any,
    out_prs: Any,
    title: str,
    body_items: list[str],
    slide_number: int,
    assets_root: Path,
) -> bool:
    """Clone template slide 5 chrome + merge the Bookmark Card infographic fragment.

    Returns True when the infographic was applied.  Returns False (without modifying
    *out_prs*) when *body_items* exceeds the 5-slot fragment capacity — the caller
    must then fall back to clone_body_block_slide so no item is dropped.

    Provides body-block chrome (title bar, logo, footer, slide number) from the
    template, then merges the 5-card Bookmark Card fragment via XML shape-tree copy
    (editable shapes — no rasterization).  Post-merge:
      - fills 5 Headline slots with *body_items* in card order (left → right)
      - clears Body lorem ipsum placeholder text
      - strips numbered oval badges (items are parallel capabilities, not steps)
      - removes the fragment's own category/title header textboxes
      - overrides fading fills with uniform brand-purple (all cards equally visible)
      - converts Freeform 6 (bookmark tail) to a roundRect matching card 2-5 height

    Parameters
    ----------
    template_prs : Presentation loaded by load_template()
    out_prs      : output Presentation to append the slide to
    title        : slide title text (fills TextBox 5 on the body-block template)
    body_items   : up to 5 parallel capability labels (one per card, left → right)
    slide_number : 1-based output slide number (updates slidenum cached field)
    assets_root  : path to the monorepo assets/ directory

    Returns
    -------
    bool — True: infographic applied.  False: too many items; caller must fall back.
    """
    if len(body_items) > _BOOKMARK_CARD_MAX_ITEMS:
        logger.info(
            "clone_infographic: %d items exceeds %d-slot Bookmark Card capacity"
            " — caller must fall back to body-block",
            len(body_items),
            _BOOKMARK_CARD_MAX_ITEMS,
        )
        return False
    from app.templates.layouts import Rect, W, H  # local import avoids circular dep
    from app.templates.merge import merge_fragment

    tmpl_slide: Any = template_prs.slides[_TEMPLATE_SLIDE_IDX]
    blank_layout = out_prs.slide_layouts[6]
    new_slide: Any = out_prs.slides.add_slide(blank_layout)

    src_cSld   = tmpl_slide.element.find(qn("p:cSld"))
    src_spTree = src_cSld.find(qn("p:spTree"))
    tgt_spTree = new_slide.shapes._spTree

    # Clear blank layout placeholders
    for child in list(tgt_spTree):
        if child.tag not in _STRUCT_TAGS:
            tgt_spTree.remove(child)

    # Deep-copy body-block chrome: title bar, logo, footer, slide-number field.
    # Skip the guide box (Rounded Rectangle 1) — it is a layout marker, not rendered.
    start_id = _next_shape_id(tgt_spTree)
    counter = start_id
    for child in src_spTree:
        if child.tag in _STRUCT_TAGS:
            continue
        if _shape_name(child) == _GUIDE_BOX_NAME:
            continue
        copied = deepcopy(child)
        for el in copied.iter():
            if el.get("id") is not None:
                el.set("id", str(counter))
                counter += 1
        tgt_spTree.append(copied)

    _register_image_rels(tmpl_slide, new_slide, tgt_spTree)
    _replace_title(tgt_spTree, title)
    _update_slide_number(tgt_spTree, slide_number)

    # Merge the Bookmark Card fragment.  Use the full slide as the target region so
    # the fragment maps 1:1 (both the fragment and the iMocha template share the same
    # 13.33" × 7.5" canvas = 12 192 000 × 6 858 000 EMU).
    frag_path = assets_root / _BOOKMARK_CARD_FRAG
    region = Rect(left=0, top=0, width=W, height=H)
    merge_fragment(new_slide, frag_path, 0, region)

    # Post-merge: fill headlines, clear bodies, strip badges and fragment header.
    # merge_fragment always appends the group as the last spTree child.
    grp_el = tgt_spTree[-1]
    _fill_bookmark_card_group(grp_el, body_items)
    return True


def clone_hub_spoke_slide(
    template_prs: Any,
    out_prs: Any,
    title: str,
    body_items: list[str],
    slide_number: int,
    assets_root: Path,
) -> bool:
    """Clone template slide 5 chrome + merge the 4-Feature Hub-and-Spoke fragment.

    Returns True when the infographic was applied.  Returns False (without modifying
    *out_prs*) when *body_items* exceeds the 4-slot fragment capacity — the caller
    must then fall back to clone_body_block_slide.

    Post-merge cleanup:
      - Removes fragment header textboxes (TextBox 2 / 3).
      - Fills the 4 wide body bars (Rectangle 69) with *body_items* text, ordered
        left-to-right, top-to-bottom.
      - Clears the narrow "Feature" headline labels (Text Placeholder 33) that are
        too small to hold full-sentence items.
      - Sets 9 pt font + latinLnBrk="0" + normAutofit on body bars.
    """
    if len(body_items) > _HUB_SPOKE_MAX_ITEMS:
        logger.info(
            "clone_hub_spoke: %d items exceeds %d-slot Hub-and-Spoke capacity"
            " — caller must fall back to body-block",
            len(body_items),
            _HUB_SPOKE_MAX_ITEMS,
        )
        return False

    from app.templates.layouts import Rect, W, H  # local import avoids circular dep
    from app.templates.merge import merge_fragment

    tmpl_slide: Any = template_prs.slides[_TEMPLATE_SLIDE_IDX]
    blank_layout = out_prs.slide_layouts[6]
    new_slide: Any = out_prs.slides.add_slide(blank_layout)

    src_spTree = tmpl_slide.element.find(qn("p:cSld")).find(qn("p:spTree"))
    tgt_spTree = new_slide.shapes._spTree

    for child in list(tgt_spTree):
        if child.tag not in _STRUCT_TAGS:
            tgt_spTree.remove(child)

    start_id = _next_shape_id(tgt_spTree)
    counter = start_id
    for child in src_spTree:
        if child.tag in _STRUCT_TAGS:
            continue
        if _shape_name(child) == _GUIDE_BOX_NAME:
            continue
        copied = deepcopy(child)
        for el in copied.iter():
            if el.get("id") is not None:
                el.set("id", str(counter))
                counter += 1
        tgt_spTree.append(copied)

    _register_image_rels(tmpl_slide, new_slide, tgt_spTree)
    _replace_title(tgt_spTree, title)
    _update_slide_number(tgt_spTree, slide_number)

    frag_path = assets_root / _HUB_SPOKE_FRAG
    region = Rect(left=0, top=0, width=W, height=H)
    merge_fragment(new_slide, frag_path, 0, region)

    grp_el = tgt_spTree[-1]
    _fill_hub_spoke_group(grp_el, body_items)
    return True


def clone_body_block_slide(
    template_prs: Any,
    out_prs: Any,
    title: str,
    body_items: list[str],
    slide_number: int,
    icon_path: Path | None = None,
    theme: ThemePalette | None = None,
    assets_root: Path | None = None,
) -> bool:
    """Clone template slide 5 into *out_prs* and fill with content.

    Parameters
    ----------
    template_prs : Presentation loaded by load_template()
    out_prs      : output Presentation to append the cloned slide to
    title        : slide title (replaces '<Slide Title>' in TextBox 5)
    body_items   : bullet / body lines to inject inside the content box
    slide_number : 1-based output slide number (updates slidenum field)
    theme        : ThemePalette (LIGHT or DARK); defaults to LIGHT
    assets_root  : path to assets/ directory; required for dark-theme logo swap

    Returns
    -------
    bool — True when content overflows even at minimum font size and the
            slide should be flagged for review.
    """
    if theme is None:
        theme = LIGHT

    is_dark = theme is not LIGHT

    tmpl_slide: Any = template_prs.slides[_TEMPLATE_SLIDE_IDX]

    # ---- add a blank slide ----
    blank_layout = out_prs.slide_layouts[6]
    new_slide: Any = out_prs.slides.add_slide(blank_layout)

    # ---- set slide background solid color (fallback / base layer) ----
    bg_fill = new_slide.background.fill
    bg_fill.solid()
    br, bg_, bb = brand.rgb_tuple(theme.bg)
    bg_fill.fore_color.rgb = RGBColor(br, bg_, bb)

    src_cSld   = tmpl_slide.element.find(qn("p:cSld"))
    src_spTree = src_cSld.find(qn("p:spTree"))
    tgt_spTree = new_slide.shapes._spTree

    # Clear any shapes the blank layout inserted (keep structural nvGrpSpPr / grpSpPr)
    for child in list(tgt_spTree):
        if child.tag not in _STRUCT_TAGS:
            tgt_spTree.remove(child)

    start_id = _next_shape_id(tgt_spTree)
    counter  = start_id

    if not is_dark:
        # Light theme — render the content panel first (lowest z-order) so chrome
        # shapes (title, logo, footer) stack on top.
        # Dark theme: the solid bg fill set above covers the slide full-bleed; no panel.
        for child in src_spTree:
            if child.tag not in _STRUCT_TAGS and _shape_name(child) == _GUIDE_BOX_NAME:
                panel_copy = deepcopy(child)
                for el in panel_copy.iter():
                    if el.get("id") is not None:
                        el.set("id", str(counter))
                        counter += 1
                _set_shape_solid_fill(panel_copy, theme.panel_fill)
                tgt_spTree.append(panel_copy)
                break

    # ---- copy all chrome shapes (excluding the panel, handled above) ----
    for child in src_spTree:
        if child.tag in _STRUCT_TAGS:
            continue
        if _shape_name(child) == _GUIDE_BOX_NAME:
            continue
        copied = deepcopy(child)
        for el in copied.iter():
            if el.get("id") is not None:
                el.set("id", str(counter))
                counter += 1
        tgt_spTree.append(copied)

    # ---- register image relationships (logo, etc.) ----
    _register_image_rels(tmpl_slide, new_slide, tgt_spTree)

    # ---- dark theme: swap template logo → logo_white.png ----
    if is_dark and assets_root is not None:
        _swap_logo_dark(new_slide, tgt_spTree, assets_root)

    # ---- fill title (preserves run formatting from template) ----
    if not _replace_title(tgt_spTree, title):
        logger.warning("clone_body_block: TextBox 5 not found — title not set")

    # ---- update cached page-number display ----
    _update_slide_number(tgt_spTree, slide_number)

    # ---- inject body textbox with overflow handling ----
    flagged = _inject_body(new_slide, body_items, theme)

    # ---- place icon accent in upper-right of content panel ----
    if icon_path is not None:
        _inject_icon(new_slide, icon_path, theme)

    return flagged


def clone_cover_slide(
    template_prs: Any,
    out_prs: Any,
    title: str,
    body_items: list[str],
) -> None:
    """Clone template slide 1 into *out_prs* and fill cover text slots.

    All template shapes (background photo, overlay graphic, logo) are preserved
    via spTree deep-copy.  Only the two placeholder text boxes are written:

    - TextBox 7  ("Cover Slide Heading Text")  → *title*
    - TextBox 9  ("Cover Slide subheading …")  → *body_items* as paragraphs
      body_items[0] = subheading line; body_items[1:] = description bullets
    - TextBox 12 ("Deck Name")                 → left as-is (editable placeholder)

    Parameters
    ----------
    template_prs : Presentation loaded by load_template()
    out_prs      : output Presentation to prepend the cover slide to
    title        : source deck's first-slide title
    body_items   : source deck's first-slide body text (subtitle + bullets)
    """
    new_slide = _clone_template_slide(template_prs, out_prs, _TEMPLATE_COVER_IDX)
    tgt_spTree = new_slide.shapes._spTree

    if not _replace_cover_heading(tgt_spTree, title):
        logger.warning("clone_cover: %r not found — heading not set", _COVER_HEADING_SHAPE)

    # TextBox 9 receives ONLY the subtitle (body_items[0]); no overflow into heading.
    if not _replace_cover_body(tgt_spTree, body_items[:1]):
        logger.warning("clone_cover: %r not found — body not set", _COVER_BODY_SHAPE)

    # Actual bullet points go in a separate textbox below the subtitle, left half only.
    if len(body_items) > 1:
        _inject_cover_bullets(new_slide, body_items[1:])


def clone_closing_slide(
    template_prs: Any,
    out_prs: Any,
) -> None:
    """Clone template slide 20 ("Skills Visibility. Business Agility.") into *out_prs*.

    The slide is appended as a fixed brand end-card with no content replacement.
    Call this AFTER all content slides have been added so it lands last.

    Parameters
    ----------
    template_prs : Presentation loaded by load_template()
    out_prs      : output Presentation to append the closing slide to
    """
    _clone_template_slide(template_prs, out_prs, _TEMPLATE_CLOSING_IDX)
