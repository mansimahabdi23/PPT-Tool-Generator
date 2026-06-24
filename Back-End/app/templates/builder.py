"""Slide builder — paints a LayoutSpec onto a blank python-pptx slide.

Public API
----------
SlideContent  -- dataclass carrying per-slide text
add_slide(prs, slide_type, content, assets_root) -> slide
"""

from __future__ import annotations

from dataclasses import dataclass, field
from pathlib import Path
from typing import Any

from lxml import etree  # type: ignore[import-untyped]
from pptx.dml.color import RGBColor
from pptx.util import Emu, Pt

from app import brand
from app.models.enums import SlideType
from app.templates.layouts import (
    LAYOUTS,
    GradientBg,
    LayoutSpec,
    Rect,
    SolidBg,
    TextSpec,
)

# OOXML namespaces used for raw XML construction
_P = "http://schemas.openxmlformats.org/presentationml/2006/main"
_A = "http://schemas.openxmlformats.org/drawingml/2006/main"


# ---------------------------------------------------------------------------
# Content descriptor
# ---------------------------------------------------------------------------

@dataclass
class SlideContent:
    """Per-slide text payload passed to add_slide."""

    title: str = ""
    subtitle: str = ""
    body_items: list[str] = field(default_factory=list)
    caption: str = ""
    kicker: str = ""


# ---------------------------------------------------------------------------
# Background helpers
# ---------------------------------------------------------------------------

def _apply_solid_bg(slide: Any, color_hex: str) -> None:
    bg = slide.background
    fill = bg.fill
    fill.solid()
    r, g, b = brand.rgb_tuple(color_hex)
    fill.fore_color.rgb = RGBColor(r, g, b)  # type: ignore[no-untyped-call]


def _apply_gradient_bg(slide: Any, spec: GradientBg) -> None:
    """Set slide background to a two-stop linear gradient via raw XML."""
    from pptx.oxml.ns import qn

    cSld: Any = slide.element.find(qn("p:cSld"))
    for existing in cSld.findall(qn("p:bg")):
        cSld.remove(existing)

    bg_xml = (
        f'<p:bg xmlns:p="{_P}" xmlns:a="{_A}">'
        f"  <p:bgPr>"
        f'    <a:gradFill rot="1">'
        f"      <a:gsLst>"
        f'        <a:gs pos="0"><a:srgbClr val="{spec.stop1_hex}"/></a:gs>'
        f'        <a:gs pos="100000"><a:srgbClr val="{spec.stop2_hex}"/></a:gs>'
        f"      </a:gsLst>"
        f'      <a:lin ang="{spec.angle}" scaled="0"/>'
        f"    </a:gradFill>"
        f"    <a:effectLst/>"
        f"  </p:bgPr>"
        f"</p:bg>"
    )
    bg_el = etree.fromstring(bg_xml)
    cSld.insert(0, bg_el)


def _apply_image_bg(slide: Any, assets_root: Path, rel_path: str) -> None:
    """Add a full-bleed picture as the background (sent to back of Z-order)."""
    img_path = assets_root / rel_path
    pic = slide.shapes.add_picture(
        str(img_path),
        Emu(0),
        Emu(0),
        width=Emu(brand.SLIDE_WIDTH_EMU),
        height=Emu(brand.SLIDE_HEIGHT_EMU),
    )
    # Move the picture element just after nvGrpSpPr (idx 0) and grpSpPr (idx 1)
    sp_tree: Any = slide.shapes._spTree
    sp_tree.remove(pic._element)
    sp_tree.insert(2, pic._element)


def _apply_background(slide: Any, spec: LayoutSpec, assets_root: Path) -> None:
    bg = spec.background
    if isinstance(bg, SolidBg):
        _apply_solid_bg(slide, bg.color_hex)
    elif isinstance(bg, GradientBg):
        _apply_gradient_bg(slide, bg)
    else:  # ImageBg
        _apply_image_bg(slide, assets_root, bg.rel_path)


# ---------------------------------------------------------------------------
# Shape helpers
# ---------------------------------------------------------------------------

def _add_textbox(slide: Any, ts: TextSpec, text: str) -> None:
    """Add a text box with a single run styled by *ts*."""
    txBox: Any = slide.shapes.add_textbox(
        Emu(ts.rect.left),
        Emu(ts.rect.top),
        Emu(ts.rect.width),
        Emu(ts.rect.height),
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.name = ts.font
    run.font.size = Pt(ts.size_pt)
    r, g, b = brand.rgb_tuple(ts.color_hex)
    run.font.color.rgb = RGBColor(r, g, b)  # type: ignore[no-untyped-call]
    run.font.bold = ts.bold


def _add_bullet_textbox(slide: Any, ts: TextSpec, items: list[str]) -> None:
    """Add a text box with one paragraph per bullet item."""
    txBox: Any = slide.shapes.add_textbox(
        Emu(ts.rect.left),
        Emu(ts.rect.top),
        Emu(ts.rect.width),
        Emu(ts.rect.height),
    )
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = para.add_run()
        run.text = item
        run.font.name = ts.font
        run.font.size = Pt(ts.size_pt)
        r, g, b = brand.rgb_tuple(ts.color_hex)
        run.font.color.rgb = RGBColor(r, g, b)  # type: ignore[no-untyped-call]


def _add_accent_bar(slide: Any, rect: Rect) -> None:
    """Add the solid brand-orange accent bar."""
    shape: Any = slide.shapes.add_shape(
        1,  # MSO_AUTO_SHAPE_TYPE.RECTANGLE
        Emu(rect.left),
        Emu(rect.top),
        Emu(rect.width),
        Emu(rect.height),
    )
    shape.fill.solid()
    r, g, b = brand.rgb_tuple(brand.BRAND_ORANGE)
    shape.fill.fore_color.rgb = RGBColor(r, g, b)  # type: ignore[no-untyped-call]
    shape.line.fill.background()  # no border


def _add_logo(slide: Any, assets_root: Path, rel_path: str, logo_rect: Rect) -> None:
    """Add a logo picture. width=0 in logo_rect means auto-scale from height."""
    logo_path = assets_root / rel_path
    if not logo_path.exists():
        return  # graceful: skip missing logo rather than crash
    w = Emu(logo_rect.width) if logo_rect.width > 0 else None
    slide.shapes.add_picture(
        str(logo_path),
        Emu(logo_rect.left),
        Emu(logo_rect.top),
        width=w,
        height=Emu(logo_rect.height),
    )


def _add_region_marker(slide: Any, region: Rect) -> None:
    """Add a dashed-border placeholder rectangle to mark the infographic region.

    Visible when no fragment has been merged — the region marker helps during
    eyeballing and is overdrawn by merge_fragment content.
    """
    shape: Any = slide.shapes.add_shape(
        1,  # RECTANGLE
        Emu(region.left),
        Emu(region.top),
        Emu(region.width),
        Emu(region.height),
    )
    shape.fill.background()  # transparent fill
    line = shape.line
    line.width = Emu(12_700)  # 1 pt
    r, g, b = brand.rgb_tuple(brand.BRAND_PURPLE)
    line.color.rgb = RGBColor(r, g, b)  # type: ignore[no-untyped-call]

    # Label
    tf = shape.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "[ INFOGRAPHIC REGION ]"
    run.font.name = brand.POPPINS
    run.font.size = Pt(10)
    r2, g2, b2 = brand.rgb_tuple(brand.BRAND_PURPLE)
    run.font.color.rgb = RGBColor(r2, g2, b2)  # type: ignore[no-untyped-call]


# ---------------------------------------------------------------------------
# Public API
# ---------------------------------------------------------------------------

def add_slide(
    prs: Any,
    slide_type: SlideType,
    content: SlideContent,
    assets_root: Path,
) -> Any:
    """Paint one iMocha-branded slide onto *prs* and return the new slide.

    Uses a blank python-pptx layout (index 6 in the default template) so no
    placeholder conflicts arise. All text, shapes, and images are added explicitly.
    """
    spec: LayoutSpec = LAYOUTS[slide_type]

    # Grab the blank layout (index 6 in the default python-pptx master)
    blank_layout: Any = prs.slide_layouts[6]
    slide: Any = prs.slides.add_slide(blank_layout)

    # 1. Background
    _apply_background(slide, spec, assets_root)

    # 2. Accent bar (agenda)
    if spec.accent_bar is not None:
        _add_accent_bar(slide, spec.accent_bar)

    # 3. Title
    _add_textbox(slide, spec.title, content.title)

    # 4. Subtitle (title / closing layouts)
    if spec.subtitle is not None and content.subtitle:
        _add_textbox(slide, spec.subtitle, content.subtitle)

    # 5. Body / bullet list
    if spec.body is not None:
        items = content.body_items or ([content.subtitle] if content.subtitle else [])
        if items:
            _add_bullet_textbox(slide, spec.body, items)

    # 6. Kicker (divider)
    if spec.kicker is not None and content.kicker:
        _add_textbox(slide, spec.kicker, content.kicker)

    # 7. Caption (data)
    if spec.caption is not None and content.caption:
        _add_textbox(slide, spec.caption, content.caption)

    # 8. Infographic region marker (shows region boundary before fragment merged)
    if spec.infographic_region is not None:
        _add_region_marker(slide, spec.infographic_region)

    # 9. Logo
    if spec.logo_rel_path and spec.logo_rect is not None:
        _add_logo(slide, assets_root, spec.logo_rel_path, spec.logo_rect)

    return slide
